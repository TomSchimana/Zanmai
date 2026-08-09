#!/usr/bin/env python3
"""Design check: the part of judging a set piece that a script decides better
than a look.

A designed document fails in two different ways. Some faults are matters of
taste, and no script has an opinion worth having about them. Others are counts,
and those are exactly what an eye glosses over on the fortieth page: how many
forms one component has grown, a colour that is in no palette, a page that runs
out of content halfway down. Whoever built the piece is also the worst judge of
it, so the counts are taken out of their hands and given to this.

Usage:
    design-check.py <kit.css> [--tokens <tokens.css>] [--pdf <render.pdf>]
                    [--fill <percent>] [--dpi <n>]

The kit declares what it intends, in one comment line the script reads. CSS or
Typst, whichever medium this format is built in:

    /* zanmai-kit: quote=q- 2, table=t- 2, code=code- 1, opener=op 1 */   (CSS)
    // zanmai-kit: quote=quote- 2, table=table- 2, code=code- 1          (Typst)
    /* zanmai-sizes: 4 */   or   // zanmai-sizes: 4

Each component gives a name prefix and the number of forms the kit allows. In CSS a
form is a class, in Typst a `#let` definition.
`zanmai-sizes` is how many distinct type sizes the scale has, counted after
custom properties are resolved, so a kit that sets every size through a variable
is checked like any other.

**Every check reports the number it worked on, and a zero is a failure.** A check
that could not run must never be indistinguishable from a check that passed: no
declaration, an empty palette, no pages measured, each of those is red, not a
polite note. That is the whole reason this script exists, since the same trap
(nothing found reads as nothing wrong) is what let the faults through in the first
place.

Exit code 1 means the piece is not finished, whatever it looks like.
"""

from __future__ import annotations

import argparse
import re
import shutil
import subprocess
import sys
import tempfile
from collections import Counter
from pathlib import Path

KIT_DECL = re.compile(r"(?:/\*|//)\s*zanmai-kit:\s*(.+?)(?:\*/|$)", re.DOTALL | re.MULTILINE)
SIZES_DECL = re.compile(r"(?:/\*|//)\s*zanmai-sizes:\s*(\d+)")
TYPST_LET = re.compile(r"^\s*#let\s+([A-Za-z][\w-]*)", re.MULTILINE)
TYPST_SIZE = re.compile(r"size\s*:\s*([^,)\]\n]+)")
TYPST_BINDING = re.compile(r"^\s*#let\s+([\w-]+)\s*=\s*([^\n]+)", re.MULTILINE)
COMPONENT = re.compile(r"([A-Za-z][\w-]*)\s*=\s*([\w-]+)\s+(\d+)")
CLASS_SELECTOR = re.compile(r"\.(-?[_a-zA-Z][\w-]*)")
HEX_COLOUR = re.compile(r"#([0-9a-fA-F]{3,8})\b")
FONT_SIZE = re.compile(r"font-size\s*:\s*([^;}]+)")
CUSTOM_PROP = re.compile(r"(--[\w-]+)\s*:\s*([^;}]+)")
VAR_REF = re.compile(r"var\(\s*(--[\w-]+)\s*(?:,[^)]*)?\)")
IMPORT_RULE = re.compile(r"""@import\s+(?:url\(\s*)?['"]([^'"]+)['"]""")


def strip_comments(css: str) -> str:
    return re.sub(r"/\*.*?\*/", " ", css, flags=re.DOTALL)


def parse_declaration(css: str) -> tuple[dict[str, tuple[str, int]], int | None]:
    """Read the kit's own statement of intent out of the comment header."""
    components: dict[str, tuple[str, int]] = {}
    match = KIT_DECL.search(css)
    if match:
        for name, prefix, allowed in COMPONENT.findall(match.group(1)):
            components[name] = (prefix, int(allowed))
    sizes = SIZES_DECL.search(css)
    return components, int(sizes.group(1)) if sizes else None


def rules(css: str) -> list[tuple[str, str]]:
    """Selector and body per rule block. Innermost braces match first, so a rule
    inside `@media` is found the same way as a top-level one, and an at-rule
    prelude simply carries no class. Line breaks and several rules on one line
    both stop mattering, which a line-based reading gets wrong."""
    return [
        (selector.strip(), body)
        for selector, body in re.findall(r"([^{}]+)\{([^{}]*)\}", strip_comments(css))
    ]


def classes_in(css: str, typst: bool = False) -> list[str]:
    """The names a form can have, in source order, without duplicates. In CSS that
    is every class in a selector; in Typst a form is a `#let` definition, which is
    the same thing one layer up: one definition per form, so a second variant shows
    up as a second name rather than hiding inside a rule."""
    seen: list[str] = []
    if typst:
        for name in TYPST_LET.findall(css):
            if name not in seen:
                seen.append(name)
        return seen
    for selector, _body in rules(css):
        for name in CLASS_SELECTOR.findall(selector):
            if name not in seen:
                seen.append(name)
    return seen


def forms_of(prefix: str, names: list[str]) -> list[str]:
    """The forms of one component, which is not simply everything starting with
    the prefix. `.q-panel-label` is a part of a quote, not a second kind of quote,
    and counting it spends the budget on an inner element. So a form is the bare
    prefix itself or the prefix plus exactly one more segment; anything with a
    further dash is a part of a form. This also keeps a prefix without a trailing
    dash from swallowing unrelated classes."""
    stem = prefix.rstrip("-")
    forms = []
    for name in names:
        if name == stem:
            forms.append(name)
            continue
        if not name.startswith(stem + "-"):
            continue
        if "-" not in name[len(stem) + 1:]:
            forms.append(name)
    return forms


def check_form_ceiling(css: str, components: dict[str, tuple[str, int]], typst: bool = False) -> tuple[list[str], int]:
    """One component, one handful of forms. Repetition is what a reader reads
    structure by, so a component that has quietly grown to five variants has
    stopped being a form at all."""
    failures = []
    names = classes_in(css, typst)
    counted = 0
    for component, (prefix, allowed) in sorted(components.items()):
        forms = forms_of(prefix, names)
        counted += len(forms)
        if not forms:
            failures.append(
                f"{component}: the kit declares prefix '{prefix}' and no class uses it, "
                "so nothing was counted for this component"
            )
        if len(forms) > allowed:
            failures.append(
                f"{component}: {len(forms)} forms, kit allows {allowed} "
                f"({', '.join('.' + f for f in sorted(forms))})"
            )
    return failures, counted


def read_css_with_imports(path: Path, depth: int = 3, seen: set[Path] | None = None) -> str:
    """A token file's own text plus everything it imports, resolved relative to
    it. The entry point a brand documents is often nothing but `@import` lines, so
    reading it literally yields an empty palette and every colour then counts as
    foreign, or a kit with no literal colours passes for the wrong reason."""
    seen = seen if seen is not None else set()
    path = path.resolve()
    if path in seen or depth < 0 or not path.is_file():
        return ""
    seen.add(path)
    try:
        text = path.read_text(encoding="utf-8")
    except OSError:
        return ""
    parts = [text]
    for target in IMPORT_RULE.findall(text):
        parts.append(read_css_with_imports(path.parent / target, depth - 1, seen))
    return "\n".join(parts)


def check_colours(css: str, tokens_css: str) -> tuple[list[str], int]:
    """A colour the palette does not contain is either a mistake or a decision
    nobody made. Both are worth stopping for."""
    allowed = {c.lower() for c in HEX_COLOUR.findall(tokens_css)}
    allowed |= {c * 2 for c in allowed if len(c) == 3}
    if not allowed:
        return ([
            "the token file holds no colour value, so there is no palette to check against. "
            "Point --tokens at the file that carries the values (following its imports if it "
            "only forwards to them)."
        ], 0)
    used = Counter(c.lower() for c in HEX_COLOUR.findall(strip_comments(css)))
    stray = sorted(c for c in used if c not in allowed and (c if len(c) != 3 else c * 2) not in allowed)
    return (
        [f"#{c} is in neither the brand tokens nor an allowed value ({used[c]}x)" for c in stray],
        len(allowed),
    )


def resolve_value(value: str, props: dict[str, str], depth: int = 5) -> str:
    """A declared value with its custom properties substituted, so a size set
    through a variable counts as the size it actually is."""
    value = value.strip()
    for _ in range(depth):
        match = VAR_REF.search(value)
        if not match:
            break
        name = match.group(1)
        if name not in props:
            break
        value = (value[: match.start()] + props[name].strip() + value[match.end():]).strip()
    return value


def check_sizes(css: str, declared: int, typst: bool = False) -> tuple[list[str], int]:
    """Four one-off sizes with a class name each are not a type scale. Counting
    only literal values would let the cleanest kit through unchecked, since it
    sets every size through a variable, so the variables are resolved first."""
    if typst:
        bindings = {name: value.strip() for name, value in TYPST_BINDING.findall(css)}
        sizes = set()
        for raw in TYPST_SIZE.findall(css):
            value = raw.strip()
            sizes.add(bindings.get(value, value))
        sizes = {s for s in sizes if re.search(r"\d", s)}
    else:
        body = strip_comments(css)
        props = {name: value for name, value in CUSTOM_PROP.findall(body)}
        sizes = {resolve_value(v, props) for v in FONT_SIZE.findall(body)}
        sizes = {s for s in sizes if s and "var(" not in s}
    if not sizes:
        return (["no font-size declaration found, so the type scale was not checked"], 0)
    if len(sizes) > declared:
        return (
            [f"{len(sizes)} distinct type sizes in use, scale declares {declared} "
             f"({', '.join(sorted(sizes))})"],
            len(sizes),
        )
    return ([], len(sizes))


def check_break_guards(css: str, components: dict[str, tuple[str, int]]) -> list[str]:
    """A quote or table that starts in one column and ends in the next reads as
    two broken things. In a multi-column layout the guard is not optional."""
    parsed = rules(css)
    if not any("column-count" in body or "columns" in body for _sel, body in parsed):
        return []
    failures = []
    names = classes_in(css)
    for component, (prefix, _allowed) in sorted(components.items()):
        for name in forms_of(prefix, names):
            bodies = [
                body for selector, body in parsed
                if name in CLASS_SELECTOR.findall(selector)
            ]
            if bodies and not any("break-inside" in b for b in bodies):
                failures.append(f".{name} ({component}) has no break-inside guard")
    return failures


def installed_font_families() -> set[str]:
    """Font family names this machine plausibly has, taken from the filenames in
    the standard font folders. Filename matching is coarse, but it needs no
    fontconfig and no extra library, and it answers the only question here: is a
    face the deck asks for present at all."""
    roots = [
        Path("/System/Library/Fonts"), Path("/Library/Fonts"), Path.home() / "Library/Fonts",
        Path("/usr/share/fonts"), Path("/usr/local/share/fonts"), Path.home() / ".fonts",
        Path.home() / ".local/share/fonts", Path("C:/Windows/Fonts"),
    ]
    names: set[str] = set()
    for root in roots:
        if not root.is_dir():
            continue
        for file in root.rglob("*"):
            if file.suffix.lower() in {".ttf", ".otf", ".ttc", ".dfont"}:
                names.add(re.split(r"[-_.]", file.stem)[0].lower())
    return names


def check_pptx(deck: Path, tokens_css: str) -> tuple[list[str], list[str]]:
    """A deck is judged by reading it, because only PowerPoint renders a .pptx the
    way PowerPoint shows it and driving the running application is the user's
    screen, not ours (operating-principles section 11). So no picture is taken:
    these are the questions the file itself answers, and each one is a count."""
    try:
        from pptx import Presentation
        from pptx.util import Emu  # noqa: F401
    except ImportError:
        return (["deck not verified: python-pptx is not installed, so nothing about this file was checked"], [])

    try:
        deck_file = Presentation(str(deck))
    except Exception as exc:  # noqa: BLE001
        return ([f"deck not readable: {type(exc).__name__}: {exc}"], [])

    allowed = {c.lower() for c in HEX_COLOUR.findall(tokens_css)}
    allowed |= {c * 2 for c in allowed if len(c) == 3}

    failures: list[str] = []
    faces_used: set[str] = set()
    slides = list(deck_file.slides)
    if not slides:
        return (["the deck holds no slide"], [])

    blank_layouts = 0
    empty_placeholders = 0
    overrides: set[str] = set()
    stray_colours: Counter = Counter()

    for number, slide in enumerate(slides, start=1):
        layout = slide.slide_layout
        if layout is None or not (layout.name or "").strip() or (layout.name or "").lower().startswith("blank"):
            blank_layouts += 1
            failures.append(f"slide {number} is built on a blank rather than a master layout")
        for shape in slide.placeholders:
            if shape.has_text_frame and not shape.text_frame.text.strip():
                empty_placeholders += 1
                failures.append(f"slide {number}: placeholder '{shape.name}' is left empty")
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            # An override only means something where there is something to
            # inherit. A placeholder takes its face and size from the layout, so
            # setting them there defeats the design; a free text box has no
            # layout behind it and must state them.
            inherits = shape.is_placeholder
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    font = run.font
                    if font.name:
                        faces_used.add(font.name)
                    if font.size is not None and inherits:
                        overrides.add(f"slide {number}: placeholder text '{run.text[:24]}' sets its own size")
                    try:
                        rgb = font.color.rgb if font.color and font.color.type is not None else None
                    except (AttributeError, ValueError):
                        rgb = None
                    if rgb is not None:
                        value = str(rgb).lower()
                        if allowed and value not in allowed:
                            stray_colours[value] += 1

    for failure in sorted(overrides):
        failures.append(failure + ", which overrides the layout it inherits")
    for value, count in sorted(stray_colours.items()):
        failures.append(f"#{value} is not in the palette ({count}x in the deck)")

    # A face the machine does not have is silently swapped for another one when
    # the deck is previewed or opened, so the deck looks wrong for a reason that
    # has nothing to do with the deck. Verified: a slide asking for an uninstalled
    # face previews in a substitute without a word.
    # The master's own theme font decides what every inheriting placeholder gets.
    # If that face is missing from the machine, every such placeholder is silently
    # swapped, and pinning a face per run is then a repair rather than a fault,
    # which is why it is not reported as one. The template is what needs fixing.
    theme_faces: set[str] = set()
    for master in deck_file.slide_masters:
        try:
            theme = master.part.part_related_by(
                "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme")
            theme_faces |= set(re.findall(r'<a:latin typeface="([^"]+)"', theme.blob.decode("utf-8", "ignore")))
        except Exception:  # noqa: BLE001
            pass
    theme_faces = {f for f in theme_faces if f and not f.startswith("+")}

    installed = installed_font_families()
    for face in sorted(theme_faces):
        if re.split(r"[ -]", face)[0].lower() not in installed:
            failures.append(
                f"the master's theme font is '{face}', which is not installed here, so every placeholder "
                "that inherits it is silently swapped. Fix the theme in the template, not slide by slide"
            )
    missing = sorted(f for f in faces_used if re.split(r"[ -]", f)[0].lower() not in installed)
    for face in missing:
        failures.append(
            f"the deck asks for '{face}', which is not installed here, so it is silently swapped "
            "for another face wherever it is opened"
        )

    checked = [
        f"{len(slides)} slides read",
        f"{len(theme_faces)} theme faces",
        f"{len(faces_used)} faces referenced, {len(faces_used) - len(missing)} of them installed",
        f"{len(allowed)} colours in the palette",
        f"{len(overrides)} layout overrides",
        f"{empty_placeholders} empty placeholders",
    ]
    # Whether a picture should have been a chart is not decidable from the file:
    # a photo is a legitimate picture. The rule stays in the skill; this does not
    # claim to verify it.
    return (failures, checked)


def check_fonts(pdf: Path) -> tuple[list[str], int]:
    """Every face the document uses has to travel inside it. A PDF that only names
    a font looks right on the machine that has that font installed, which is always
    the machine it was built on, and falls back to something else in the hands of
    whoever it was made for. That makes it undeliverable, and nothing about it is
    visible to the person who built it, so it is checked here rather than trusted.
    """
    if shutil.which("pdffonts") is None:
        return (["fonts not verified: pdffonts (poppler) is not on this machine, so embedding is unproven"], 0)
    result = subprocess.run(["pdffonts", str(pdf)], capture_output=True, text=True)
    if result.returncode != 0:
        return ([f"fonts not verified: pdffonts failed ({result.stderr.strip()})"], 0)
    lines = [l for l in result.stdout.splitlines()[2:] if l.strip()]
    if not lines:
        return (["the render uses no embedded font at all, which means its text is not text"], 0)
    failures = []
    for line in lines:
        parts = line.split()
        # The trailing columns are fixed (emb sub uni object-id gen) while the type
        # and encoding columns hold a varying number of words, so the row is read
        # from the end. Counting from the front puts "CID Type 0C" out by one and
        # reports every embedded font as missing.
        if len(parts) < 6:
            failures.append(f"cannot read the font row: {line.strip()}")
            continue
        if parts[-5] != "yes":
            failures.append(f"{parts[0]} is not embedded in the render, so the file cannot be handed on")
    return (failures, len(lines))


def page_fill(pdf: Path, dpi: int, threshold: float, open_pages: set[int]) -> tuple[list[str], int, int]:
    """How far down the page the content actually reaches. Half a page of white
    below the last line is a fault, and it is the one an expert reports as a
    percentage instead of fixing. Returns the pages measured, because zero pages
    measured is a failed check, not a passed one."""
    if shutil.which("pdftoppm") is None:
        return ["page fill not measured: pdftoppm (poppler) is not on this machine"], 0, 0
    try:
        from PIL import Image
    except ImportError:
        return ["page fill not measured: Pillow is not installed"], 0, 0

    failures: list[str] = []
    measured = 0
    colour_pages = 0
    with tempfile.TemporaryDirectory() as tmp:
        stem = Path(tmp) / "page"
        result = subprocess.run(
            ["pdftoppm", "-png", "-r", str(dpi), str(pdf), str(stem)],
            capture_output=True,
            text=True,
        )
        if result.returncode != 0:
            return [f"page fill not measured: pdftoppm failed ({result.stderr.strip()})"], 0, 0
        pages = sorted(Path(tmp).glob("page*.png"))
        for index, png in enumerate(pages):
            measured += 1
            with Image.open(png) as im:
                grey = im.convert("L")
                width, height = grey.size
                pixels = grey.load()
                background = Counter(
                    pixels[x, 0] for x in range(0, width, max(1, width // 64))
                ).most_common(1)[0][0]
                # A page whose ground is a colour is a colour surface, a cover or a
                # closing statement, and coverage says nothing about it: everything
                # matches the ground, so the measure would call the fullest page in
                # the document empty.
                if background < 240:
                    colour_pages += 1
                    continue
                # Coverage over a coarse grid, not the lowest inked row. In a
                # two-column layout one full column reaches the bottom while the
                # other is white from the middle down, and a lowest-row measure
                # calls that page full. The grid sees the empty half, which is what
                # the reader sees.
                cols, rows = 4, 16
                inked = 0
                for cy in range(rows):
                    for cx in range(cols):
                        x0, x1 = cx * width // cols, (cx + 1) * width // cols
                        y0, y1 = cy * height // rows, (cy + 1) * height // rows
                        step_x = max(1, (x1 - x0) // 12)
                        step_y = max(1, (y1 - y0) // 12)
                        if any(
                            abs(pixels[x, y] - background) > 12
                            for y in range(y0, y1, step_y)
                            for x in range(x0, x1, step_x)
                        ):
                            inked += 1
                fill = 100.0 * inked / (cols * rows)
                number = int(re.sub(r"\D", "", png.stem) or 0)
                # The last page of a document runs out where the text runs out, so
                # it is partial by nature and flagging it would train the reader of
                # this output to ignore the list.
                if index == len(pages) - 1:
                    continue
                if fill < threshold and number not in open_pages:
                    failures.append(f"{png.stem}: content covers {fill:.0f}% of the page")
    if not measured:
        failures.append("page fill not measured: the render produced no page images")
    return failures, measured, colour_pages


def main(argv: list[str]) -> int:
    ap = argparse.ArgumentParser(description="Check a design kit and its render against what the kit itself declares.")
    ap.add_argument("kit", type=Path, nargs="?",
                    help="the kit for this brand and format, the CSS or the Typst source; a .typ file is read as Typst")
    ap.add_argument("--tokens", type=Path, required=True,
                    help="the brand token CSS, so stray colours can be found; its @import lines are followed")
    ap.add_argument("--pdf", type=Path,
                    help="the rendered PDF, so page coverage is measured on real pages")
    ap.add_argument("--pptx", type=Path,
                    help="a native deck, checked by reading the file; no render, no application")
    ap.add_argument("--fill", type=float, default=70.0, help="minimum page coverage in percent (default 70)")
    ap.add_argument("--open-pages", default="",
                    help="Page numbers that are meant to be sparse, typed by hand. A last resort: whoever is being checked decides what is not checked, so a green result says less. Prefer --open-pages-from.")
    ap.add_argument("--open-pages-from", dest="open_pages_from",
                    help="File the build wrote from marks in the document source, one page number per line or a JSON array. The intent then lives in the source, where it is visible and reviewable, rather than in this call.")
    ap.add_argument("--dpi", type=int, default=50, help="render resolution for the fill measurement (default 50)")
    args = ap.parse_args(argv[1:])

    if args.pptx:
        if not args.pptx.is_file():
            print(f"design-check: no deck at {args.pptx}", file=sys.stderr)
            return 2
        if not args.tokens.is_file():
            print(f"design-check: no token CSS at {args.tokens}", file=sys.stderr)
            return 2
        deck_failures, deck_checked = check_pptx(args.pptx, read_css_with_imports(args.tokens))
        print("checked: " + "; ".join(deck_checked) if deck_checked else "checked: nothing")
        if deck_failures:
            print()
            for failure in deck_failures:
                print(f"FAIL: {failure}")
            print(f"\n{len(deck_failures)} open, the deck is not finished.")
            return 1
        print("\nevery point above passes.")
        return 0

    if not args.kit or not args.kit.is_file():
        print(f"design-check: no kit at {args.kit}", file=sys.stderr)
        return 2
    if not args.pdf:
        print("design-check: --pdf is required when checking a kit and its render", file=sys.stderr)
        return 2
    if not args.tokens.is_file():
        print(f"design-check: no token CSS at {args.tokens}", file=sys.stderr)
        return 2
    if not args.pdf.is_file():
        print(f"design-check: no render at {args.pdf}", file=sys.stderr)
        return 2

    css = args.kit.read_text(encoding="utf-8")
    typst = args.kit.suffix.lower() == ".typ"
    tokens = read_css_with_imports(args.tokens)
    components, declared_sizes = parse_declaration(css)

    failures: list[str] = []
    checked: list[str] = []

    if not components:
        failures.append(
            "no zanmai-kit declaration in this CSS, so no form ceiling exists to hold the build to. "
            "Add one line: /* zanmai-kit: quote=q- 2, table=t- 2, code=code- 1, opener=op 1 */"
        )
    else:
        ceiling_failures, forms_counted = check_form_ceiling(css, components, typst)
        failures += ceiling_failures
        # A break guard is a CSS problem. In a typesetting source the renderer
        # handles a block that does not fit, so there is nothing to declare.
        if not typst:
            failures += check_break_guards(css, components)
        checked.append(f"{forms_counted} forms across {len(components)} components")

    colour_failures, palette = check_colours(css, tokens)
    failures += colour_failures
    checked.append(f"{palette} colours in the palette")

    if declared_sizes is None:
        failures.append(
            "no zanmai-sizes declaration, so the type scale was not checked. "
            "Add one line: /* zanmai-sizes: 4 */"
        )
    else:
        size_failures, sizes_found = check_sizes(css, declared_sizes, typst)
        failures += size_failures
        checked.append(f"{sizes_found} type sizes in use against {declared_sizes} declared")

    font_failures, fonts = check_fonts(args.pdf)
    failures += font_failures
    checked.append(f"{fonts} fonts checked for embedding")

    open_pages = {int(p) for p in re.findall(r"\d+", args.open_pages)}
    by_hand = len(open_pages)
    derived = 0
    if args.open_pages_from:
        marks = Path(args.open_pages_from)
        if not marks.is_file():
            failures.append(
                f"--open-pages-from points at {marks}, which does not exist, so no page was "
                "exempted from the coverage check by anything the build measured"
            )
        else:
            found = {int(n) for n in re.findall(r"\d+", marks.read_text(encoding="utf-8"))}
            derived = len(found)
            open_pages |= found
    fill_failures, pages, colour_pages = page_fill(args.pdf, args.dpi, args.fill, open_pages)
    failures += fill_failures
    checked.append(
        f"{pages} pages measured for coverage at {args.dpi} dpi"
        + (f", {colour_pages} colour surfaces" if colour_pages else "")
        + (f", {derived} open by marks in the source" if derived else "")
        + (f", {by_hand} excluded by hand" if by_hand else "")
    )

    print("checked: " + "; ".join(checked))
    if failures:
        print()
        for failure in failures:
            print(f"FAIL: {failure}")
        print(f"\n{len(failures)} open, the piece is not finished.")
        return 1
    print("\nevery point above passes.")
    return 0


if __name__ == "__main__":
    sys.exit(main(sys.argv))
