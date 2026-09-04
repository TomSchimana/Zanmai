#!/usr/bin/env python3
"""A brand's own slides as a library: harvest what exists, then fill copies of it.

The expensive way to make a deck is to compose every slide from scratch, and it
is also the way that drifts: the fourth slide invents what the first one already
solved. The cheap way is the one a designer uses, which is to take the slide that
already carries this shape of content and put the new content in it. Nothing is
drawn, so nothing can drift, and a slide costs seconds.

What this does not do is ship a fixed set of layouts. A brand's look is not ours
to prescribe, so the library is read out of the material the user already has,
their template and their approved decks, and it grows as they approve more.

    slide-library.py harvest <deck.pptx> --into <library-dir>
    slide-library.py build <plan.json> <out.pptx> --library <library-dir>
    slide-library.py show <library-dir> [--slide <id>]
    slide-library.py check <library-dir> --task <slug>
    slide-library.py nudge <deck.pptx> --shape <name> --dx <in> --dy <in> --into <out.pptx>
    slide-library.py overlap-check <deck.pptx> [--slide <n>]
    slide-library.py align-check <deck.pptx> [--slide <n>]
    slide-library.py layout-check <deck.pptx> [--slide <n>]
    slide-library.py migrate <deck.pptx> --slide <n> --into <brand.pptx> --out <new.pptx>
    slide-library.py render <deck.pptx> --into <dir>

`harvest` writes `index.json` and a readable `index.md` describing each slide:
which master layout it uses, what its slots are, and how much text each slot
actually holds, measured from the box and the type size that is in it rather than
declared by anyone. That measurement is what lets a later build refuse content
that would overflow, instead of producing a slide nobody looks at twice.

`build` takes a plan naming a source slide per target slide plus the text per
slot, clones the source and swaps the text. `show` prints the library so an agent
can choose. `check` does the same and additionally records that the library was
looked at for a given piece of work (`--task`, the `workbench/<slug>/` bundle this
deliverable belongs to): `library-check-guard` (PreToolUse Bash, see
`zanmai.py hook`) refuses to save a `.pptx` under that bundle until this has run
at least once, so composing from scratch is a choice made after looking, not
instead of it.

`nudge` and `overlap-check` are for the other kind of change: a correction, not a build. Found
2026-08-25 on a real Carol run: a text block sat over its card's background numeral, a fix any
person makes by dragging the box in PowerPoint in ten seconds, and it cost fifteen minutes because
the group-coordinate math, the "does the glyph paint past its box" check (`wrap="none"` at a large
point size does exactly that) and the pairwise overlap check were re-derived from nothing, in a
freshly written script that then needed its own debugging. `nudge` moves one shape by a delta and
sets all four xfrm values explicitly, `overlap-check` finds text-over-text pairs the same way, ink
rather than the box where `wrap="none"` makes the box lie. Both work on shapes nested in groups: a
group's own frame is grown to keep containing what moved, the way `gruppe_nachziehen` did by hand.

`align-check` is the same "box lies about the ink" pattern found a second time the same day, at the
other edge: a title and a body text sitting on the identical box-left value can still read as
mis-aligned, because a large bold face and a small regular face carry different left side bearings.
Both checks measure the real ink where they can, loading the actual font file (matched by name in the
standard font folders) and reading its own metrics, not a guess about where letters typically start.
"""

from __future__ import annotations

import argparse
import copy
import json
import math
import re
import sys
from datetime import datetime, timezone
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.enum.text import MSO_AUTO_SIZE
    from pptx.oxml.ns import qn
    from pptx.util import Emu
except ImportError:
    sys.exit("slide-library: python-pptx is missing. Provision it first (tools ensure python_pptx).")

EMU_PER_INCH = 914400
# Characters per square inch of text box at 12 pt, calibrated against the
# reference build so its own slides sit at roughly half their capacity: a card
# title box of 2.57 x 0.96 inch at 18 pt carries about 46 characters and holds 40
# today. Coarse on purpose, it only has to tell "fits" from "will overflow", and
# it is calibrated rather than guessed because a number nobody measured would let
# every overflow through.
CHAR_DENSITY = 42.0
# Average glyph width of the face CHAR_DENSITY was calibrated against, in pixels per
# character at 100 pt, measured on a reference sentence. Arial and Helvetica both sit here.
# A face that runs wider carries fewer characters in the same box, and the density has to
# follow: Montserrat runs 15.5% wider than Arial, so a cell that held
# 23 characters before a brand's theme was applied holds 20 after it. Without this the count
# does not move at all when the typeface changes, and a check that cannot see a theme swap
# is quiet exactly where a brand hand-over goes wrong.
REFERENCE_GLYPH_WIDTH = 44.85
_WIDTH_SAMPLE = ("Lorem ipsum dolor sit amet consectetur adipiscing elit sed do eiusmod "
                 "tempor incididunt")
_width_cache: dict = {}


def inches(value) -> float:
    return round((value or 0) / EMU_PER_INCH, 3)


def text_of(shape) -> str:
    return shape.text_frame.text.strip() if shape.has_text_frame else ""


def sizes_in(shape) -> list[float]:
    sizes = []
    if not shape.has_text_frame:
        return sizes
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if run.font.size is not None:
                sizes.append(run.font.size.pt)
    return sizes


def face_width_factor(family: str | None, bold: bool = False) -> float:
    """How much wider this face runs than the one CHAR_DENSITY was calibrated on.

    1.0 when the face cannot be measured on this machine, which leaves the calibrated
    number exactly as it was rather than guessing a correction.
    """
    if not family:
        return 1.0
    key = (family.lower(), bold)
    if key in _width_cache:
        return _width_cache[key]
    faktor = 1.0
    try:
        from PIL import ImageFont
        path = _font_file(family, bold)
        if path is not None:
            font = ImageFont.truetype(str(path), 100)
            breite = font.getlength(_WIDTH_SAMPLE) / len(_WIDTH_SAMPLE)
            if breite > 0:
                faktor = breite / REFERENCE_GLYPH_WIDTH
    except Exception:  # noqa: BLE001 -- an unmeasurable face falls back, it does not crash a check
        faktor = 1.0
    _width_cache[key] = faktor
    return faktor


def capacity(shape) -> int:
    """How many characters this box carries at the size its own text uses. A slot
    with no text yet falls back to the deck's body size, which is the honest
    answer rather than a guess dressed as a number."""
    size = max(sizes_in(shape) or [12.0])
    area = inches(shape.width) * inches(shape.height)
    return int(area * CHAR_DENSITY * (12.0 / size) ** 2)


def cell_sizes(cell) -> list[float]:
    sizes = []
    for paragraph in cell.text_frame.paragraphs:
        for run in paragraph.runs:
            if run.font.size is not None:
                sizes.append(run.font.size.pt)
    return sizes


def cell_capacity(cell, width_emu, height_emu, family: str | None = None) -> int:
    """Same formula as `capacity()`, for a table cell rather than a shape: a cell has no
    `.width`/`.height` of its own, those live on the table's column and row.

    `family` is the typeface the cell actually renders in. Given, the count is corrected by
    how wide that face runs against the calibration face; left out, the calibrated number
    stands unchanged. This is what lets the check notice a theme swap: the same table in
    Montserrat holds about 13% less than in Arial, and before this the number did not move.
    """
    size = max(cell_sizes(cell) or [12.0])
    area = inches(width_emu) * inches(height_emu)
    roh = area * CHAR_DENSITY * (12.0 / size) ** 2
    return int(roh / face_width_factor(family))


def contains(outer, inner) -> bool:
    return (outer.left <= inner.left and outer.top <= inner.top
            and outer.left + outer.width >= inner.left + inner.width
            and outer.top + outer.height >= inner.top + inner.height)


def rect_contains(outer, inner) -> bool:
    """Same test on (left, top, width, height) tuples, for shapes whose own `left` is
    expressed in their group's child space rather than in slide space."""
    ol, ot, ow, oh = outer
    il, it, iw, ih = inner
    return ol <= il and ot <= it and ol + ow >= il + iw and ot + oh >= it + ih


def _every_leaf(shapes, tf=(1.0, 0.0, 1.0, 0.0)):
    """(shape, rect-in-slide-EMU) for every non-group shape, groups resolved.

    `_leaf_rects` does the same for text-bearing leaves only. Naming a slot needs the
    containers too -- the filled boxes a card is drawn as carry no text and are exactly
    what gives the text inside them its name.

    This exists because everything above worked on `slide.shapes`, which is the top level
    only. Measured on a corporate template: `slots` found 6 places where
    the deck has about 30, and one banded slide reported a single slot, its title, while
    five labelled bands sat on it. That template draws almost every rubric as a group
    inside a group, so the whole build chain -- harvest, slots, fill -- saw nothing.
    """
    ax, bx, ay, by = tf
    out = []
    for sh in shapes:
        if sh.shape_type == 6:
            sx, sy = _group_scale(sh)
            xfrm = sh._element.grpSpPr.xfrm
            off, chOff = xfrm.find(qn("a:off")), xfrm.find(qn("a:chOff"))
            nested = (ax * sx, ax * (int(off.get("x")) - int(chOff.get("x")) * sx) + bx,
                      ay * sy, ay * (int(off.get("y")) - int(chOff.get("y")) * sy) + by)
            out += _every_leaf(sh.shapes, nested)
            continue
        if sh.left is None:
            continue
        out.append((sh, (ax * sh.left + bx, ay * sh.top + by, ax * sh.width, ay * sh.height)))
    return out


def slot_names(slide) -> dict:
    """Name every text-bearing shape by the role its geometry gives it.

    The grouping is what makes a library usable: a card is a filled box with text
    boxes sitting inside it, so the boxes belong to that card and get its number,
    left to right. Within a group the largest type is the title and the rest is
    the body. No naming convention in the source file is required, which matters
    because the source is the user's own deck, not something we authored."""
    # Every shape, groups resolved into slide space. Working on `slide.shapes` alone was the
    # long-standing defect: a template that draws its rubrics as groups showed almost no slots.
    blaetter = _every_leaf(slide.shapes)
    frames = [(s, r) for s, r in blaetter if s.has_text_frame and text_of(s)]
    containers = [(s, r) for s, r in blaetter
                  if (not s.has_text_frame or not text_of(s)) and r[2] and r[3]]

    groups: dict[int, list] = {}
    loose = []
    for frame, frame_rect in frames:
        holder = None
        for index, (_container, c_rect) in enumerate(containers):
            if rect_contains(c_rect, frame_rect):
                if holder is None or (containers[holder][1][2] * containers[holder][1][3]
                                      > c_rect[2] * c_rect[3]):
                    holder = index
        if holder is None:
            loose.append((frame, frame_rect))
        else:
            groups.setdefault(holder, []).append((frame, frame_rect))

    slots: dict[str, dict] = {}

    def add(name, shape, rect):
        # A name that repeats would silently overwrite the first shape carrying it, and `fill`
        # would then write into one place and leave the other standing.
        if name in slots:
            zaehler = 2
            while f"{name}{zaehler}" in slots:
                zaehler += 1
            name = f"{name}{zaehler}"
        slots[name] = {
            "text": text_of(shape),
            "capacity": capacity(shape),
            "lines": len([p for p in shape.text_frame.paragraphs if p.text.strip()]),
            "size_pt": max(sizes_in(shape) or [0]) or None,
            "box": [inches(rect[0]), inches(rect[1]), inches(rect[2]), inches(rect[3])],
        }

    for shape, rect in sorted(loose, key=lambda sr: (sr[1][1], sr[1][0])):
        if shape.is_placeholder and shape.placeholder_format.idx == 0:
            add("title", shape, rect)
        else:
            add(f"text{sum(1 for k in slots if k.startswith('text')) + 1}", shape, rect)

    ordered = sorted(groups.items(), key=lambda kv: (containers[kv[0]][1][1], containers[kv[0]][1][0]))
    rows: dict[float, list] = {}
    for holder, members in ordered:
        rows.setdefault(round(inches(containers[holder][1][1]), 1), []).append((holder, members))
    for row_index, (top, entries) in enumerate(sorted(rows.items()), start=1):
        entries.sort(key=lambda kv: containers[kv[0]][1][0])
        prefix = "hub" if len(entries) == 1 and row_index == 1 else "card"
        for column, (holder, members) in enumerate(entries, start=1):
            name = prefix if prefix == "hub" else f"{prefix}{column}"
            members.sort(key=lambda sr: -(max(sizes_in(sr[0]) or [0])))
            for position, (member, m_rect) in enumerate(members):
                add(f"{name}.title" if position == 0 else f"{name}.lines{position if position > 1 else ''}",
                    member, m_rect)

    # A template's rubrics are often laid out as a table rather than free text boxes; a table's
    # cells carry no `has_text_frame`, so they are invisible to everything above this point. Found
    # a real Battlecard template was entirely tables, and harvest reported only a title
    # and a bar, which read as "no placeholders here" and was wrong.
    table_index = 0
    for shape in slide.shapes:
        if not getattr(shape, "has_table", False):
            continue
        table_index += 1
        table = shape.table
        row_tops = []
        top = shape.top or 0
        for row in table.rows:
            row_tops.append(top)
            top += row.height or 0
        col_lefts = []
        left = shape.left or 0
        for column in table.columns:
            col_lefts.append(left)
            left += column.width or 0
        for r in range(len(table.rows)):
            for c in range(len(table.columns)):
                cell = table.cell(r, c)
                if cell.is_spanned or not cell.text_frame.text.strip():
                    continue
                width = table.columns[c].width or 0
                height = table.rows[r].height or 0
                slots[f"table{table_index}.r{r + 1}c{c + 1}"] = {
                    "text": cell.text_frame.text.strip(),
                    "capacity": cell_capacity(cell, width, height),
                    "lines": len([p for p in cell.text_frame.paragraphs if p.text.strip()]),
                    "size_pt": max(cell_sizes(cell) or [0]) or None,
                    "box": [inches(col_lefts[c]), inches(row_tops[r]), inches(width), inches(height)],
                }
    return slots


def harvest(deck_path: Path, into: Path) -> int:
    deck = Presentation(str(deck_path))
    into.mkdir(parents=True, exist_ok=True)
    entries = []
    for number, slide in enumerate(deck.slides, start=1):
        slots = slot_names(slide)
        fills = []
        for shape in slide.shapes:
            try:
                if shape.fill.type == 1:
                    value = str(shape.fill.fore_color.rgb).lower()
                    if value not in fills:
                        fills.append(value)
            except Exception:  # noqa: BLE001
                pass
        repeated = sorted({re.sub(r"\d+", "", k.split(".")[0]) for k in slots if "." in k})
        entries.append({
            "id": number,
            "layout": slide.slide_layout.name,
            "title": slots.get("title", {}).get("text", ""),
            "shapes": len(slide.shapes),
            "groups": {kind: len({k.split(".")[0] for k in slots if k.startswith(kind)}) for kind in repeated},
            "fills": fills,
            "slots": slots,
        })

    index = {"source": str(deck_path), "slides": entries}
    (into / "index.json").write_text(json.dumps(index, indent=2, ensure_ascii=False) + "\n", encoding="utf-8")

    lines = ["# Slide library", "",
             f"Harvested from `{deck_path.name}`. Every slide here is one the user has approved, so a new",
             "piece starts by taking the one that already carries this shape of content, not by composing.",
             "The capacity per slot is measured from the box and the type size in it, so a slot that would",
             "overflow is refused at build time rather than delivered.", ""]
    for entry in entries:
        groups = ", ".join(f"{count}x {kind}" for kind, count in entry["groups"].items()) or "no repeated group"
        lines.append(f"## Slide {entry['id']}: {entry['title'] or '(no title)'}")
        lines.append(f"Layout `{entry['layout']}` · {entry['shapes']} shapes · {groups}")
        lines.append("")
        lines.append("| slot | holds | now | at |")
        lines.append("|---|---|---|---|")
        for name, slot in entry["slots"].items():
            lines.append(f"| `{name}` | {slot['capacity']} chars | {len(slot['text'])} chars, "
                         f"{slot['lines']} lines | {slot['size_pt'] or '-'} pt |")
        lines.append("")
    (into / "index.md").write_text("\n".join(lines), encoding="utf-8")

    print(f"harvested {len(entries)} slides into {into}")
    for entry in entries:
        groups = ", ".join(f"{count}x {kind}" for kind, count in entry["groups"].items()) or "-"
        print(f"  {entry['id']:>3}  {entry['title'][:44]:<44} {groups}")
    return 0


_REL_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
_REL_ATTRS = tuple(f"{{{_REL_NS}}}{name}" for name in ("embed", "link", "id", "pict", "dm", "lo", "qs", "cs"))


def _carry_relationships(new_slide, source_slide, element) -> None:
    """Bring the parts a copied shape points at across, and renumber the ids to match.

    A picture is two things: the shape XML on the slide, and the image part it points at through a
    relationship id. Deep-copying the element brings the first and leaves the second behind, so the
    copy carries `r:embed="rId2"` into a slide whose rels have no rId2. PowerPoint calls that file
    damaged and strips the shape; LibreOffice renders it silently without the image, which is why
    every check and every render looked clean. Found on a finished deck, with
    the logo missing and the file repaired on open.
    """
    for el in element.iter():
        for attr in _REL_ATTRS:
            alt_id = el.get(attr)
            if not alt_id:
                continue
            try:
                rel = source_slide.part.rels[alt_id]
            except KeyError:
                continue
            if rel.is_external:
                neu_id = new_slide.part.relate_to(rel.target_ref, rel.reltype, is_external=True)
            else:
                neu_id = new_slide.part.relate_to(rel.target_part, rel.reltype)
            el.set(attr, neu_id)


def clone(deck, source):
    """A copy of one slide, shape for shape. Same layout, then every element of
    the source deep-copied in, so fills, connectors, pills and geometry come along
    exactly. Nothing is redrawn, which is the whole point.

    The relationships come with it. Copying the XML alone leaves every picture, hyperlink and
    embedded object pointing at an id the new slide does not have."""
    new = deck.slides.add_slide(source.slide_layout)
    for shape in list(new.shapes):
        shape._element.getparent().remove(shape._element)
    for shape in source.shapes:
        kopie = copy.deepcopy(shape._element)
        _carry_relationships(new, source, kopie)
        new.shapes._spTree.append(kopie)
    return new


RENDER_HINT = ("install LibreOffice: macOS `brew install --cask libreoffice`, "
               "Windows `winget install TheDocumentFoundation.LibreOffice`, "
               "Debian/Ubuntu `apt install libreoffice-impress`")


def _stand_datei(deck: Path) -> Path | None:
    """Where the note about a deck's last tool-written state lives, or nothing outside a space."""
    for oben in [deck.resolve()] + list(deck.resolve().parents):
        if (oben / "zanmai" / "system").is_dir():
            ziel = oben / "zanmai" / "runtime" / "decks.json"
            ziel.parent.mkdir(parents=True, exist_ok=True)
            return ziel
    return None


def _stand(deck: Path) -> str:
    """The file's content, not its size and timestamp: two saves in the same second are the normal
    case here, and both of those stay identical across one."""
    import hashlib
    return hashlib.md5(deck.read_bytes()).hexdigest()


def guarded_save(deck_obj, ziel: Path) -> bool:
    """Write, unless somebody else wrote to this file since a tool last did.

    Two runs on one deck is the normal case here: an expert holds it open for minutes while the
    user has a single slide corrected beside it. Nothing noticed that, and the later save took the
    earlier correction away in silence. So every tool write leaves a note of what it left behind;
    if the file no longer matches that note, somebody edited it in between and this save would
    erase their work. It reports instead, and the caller decides.

    A deck outside a space has nowhere to keep the note, and is written as before.
    """
    notiz = _stand_datei(ziel)
    vorher = None
    if notiz is not None and ziel.exists():
        try:
            alle = json.loads(notiz.read_text(encoding="utf-8"))
            vorher = alle.get(str(ziel.resolve()))
        except (OSError, json.JSONDecodeError):
            alle = {}
        if vorher and vorher != _stand(ziel):
            alle[str(ziel.resolve())] = _stand(ziel)
            notiz.write_text(json.dumps(alle, indent=1), encoding="utf-8")
            print(f"slide-library: {ziel.name} changed since a tool last wrote it, so saving now "
                  f"would take that change away. Nothing was written. Look at the file; running "
                  f"the same command again writes, because the warning has been given.",
                  file=sys.stderr)
            return False
    ziel.parent.mkdir(parents=True, exist_ok=True)
    deck_obj.save(str(ziel))
    if notiz is not None:
        try:
            alle = json.loads(notiz.read_text(encoding="utf-8")) if notiz.exists() else {}
        except (OSError, json.JSONDecodeError):
            alle = {}
        alle[str(ziel.resolve())] = _stand(ziel)
        notiz.write_text(json.dumps(alle, indent=1), encoding="utf-8")
    return True


# The order `a:spPr` declares its children in (CT_ShapeProperties). PowerPoint holds a file to it,
# LibreOffice does not, which is what made this class of fault invisible to every render.
_SPPR_ORDER = ["xfrm", "custGeom", "prstGeom", "noFill", "solidFill", "gradFill", "blipFill",
               "pattFill", "grpFill", "ln", "effectLst", "effectDag", "scene3d", "sp3d", "extLst"]


def _shape_props(shape):
    for tag in (qn("p:spPr"), qn("a:spPr")):
        gefunden = shape._element.find(tag)
        if gefunden is not None:
            return gefunden
    return None


def schema_check(deck_path: Path, slide_no: int | None = None) -> int:
    """Shapes PowerPoint will not draw, though every render shows them.

    Found in practice: on six slides of a finished deck the load-bearing shapes, a hub
    circle, a chevron chain, every bullet marker, were invisible in PowerPoint and complete in the
    file. The geometry had ended up before the position inside `a:spPr`, and against the schema
    PowerPoint drops the position and draws nothing. LibreOffice is tolerant there, so the control
    render looked right while the customer saw empty slides, and none of the five other checks here
    reads the element order. This one reads the file rather than a picture of it.
    """
    from lxml import etree

    deck = Presentation(str(deck_path))
    probleme: list[str] = []
    geprueft = 0
    for i, folie in enumerate(deck.slides, 1):
        if slide_no and i != slide_no:
            continue
        for shape, _rect in _every_leaf(folie.shapes):
            props = _shape_props(shape)
            if props is None:
                continue
            geprueft += 1
            kinder = [etree.QName(kind).localname for kind in props]
            if not shape.is_placeholder and "xfrm" not in kinder:
                probleme.append(f"slide {i}: {shape.name} carries no position (a:xfrm), so "
                                f"PowerPoint does not draw it")
                continue
            rang = [_SPPR_ORDER.index(k) for k in kinder if k in _SPPR_ORDER]
            if rang != sorted(rang):
                probleme.append(f"slide {i}: {shape.name} has a:spPr out of schema order "
                                f"({', '.join(kinder)}), so PowerPoint drops its position")
    for zeile in probleme:
        print(f"FAIL: {zeile}")
    print(f"{geprueft} shape(s) checked on {len(deck.slides.__iter__.__self__._sldIdLst)} slide(s), "
          f"{len(probleme)} that PowerPoint would not draw")
    return 1 if probleme else 0


def _nichts_geprueft(anzahl: int) -> str:
    """The sentence a check adds when it compared nothing at all.

    Zero faults out of zero comparisons is not a pass and reads exactly like one. It happened on an
    approved piece: a correction had shifted every text frame by its own side bearing, so no two
    frames shared an edge any more, the alignment check's own precondition never fired, and the
    line said `0 groups checked, 0 misalignments`. The file was fine and the check was right, and
    the line still told the reader something it had not established.
    """
    return "" if anzahl else "  NOTE: nothing matched this check's precondition, so this is not a pass, it is silence."


def _soffice() -> str | None:
    """LibreOffice on this machine, whatever it is called here."""
    import shutil
    for name in ("soffice", "libreoffice"):
        pfad = shutil.which(name)
        if pfad:
            return pfad
    for pfad in ("/Applications/LibreOffice.app/Contents/MacOS/soffice",
                 r"C:\Program Files\LibreOffice\program\soffice.exe",
                 r"C:\Program Files (x86)\LibreOffice\program\soffice.exe"):
        if Path(pfad).is_file():
            return pfad
    return None


# Where a machine keeps its typefaces, and where a run may add its own. LibreOffice on macOS reads
# fontconfig, and fontconfig's own defaults there do not include the per-user folder: a deck set in
# a face installed for the user alone came out of a headless convert in a serif nobody chose, while
# the same file opened in PowerPoint was correct. Measured on this machine: the PDF carried
# `LinuxLibertineG` where it should have carried the brand face.
#
# The fix is a configuration file of our own, pointed at with `FONTCONFIG_FILE`, listing the
# folders that actually hold fonts here plus any the caller names. Nothing is installed and nothing
# on the machine changes, which matters: the alternative that was being used instead is driving the
# real application by remote control, and that throws a window open on the screen of whoever is
# working there. A tool that interrupts the person at the keyboard is a last resort, not a method.
_FONT_DIRS = ("~/Library/Fonts", "/Library/Fonts", "/System/Library/Fonts",
              "~/.fonts", "~/.local/share/fonts", "/usr/share/fonts", "/usr/local/share/fonts",
              "C:/Windows/Fonts")


def _fontconfig_for(extra_dirs: list[Path] | None, workdir: Path) -> dict:
    """Environment for a headless office run that can see the fonts this piece is set in.

    Returns the variables to add. Extra folders come first, so a brand's own copy of a face wins
    over an older one that happens to be installed.
    """
    dirs: list[str] = []
    for d in (extra_dirs or []):
        if d and Path(d).is_dir():
            dirs.append(str(Path(d).resolve()))
    for d in _FONT_DIRS:
        pfad = Path(d).expanduser()
        if pfad.is_dir():
            dirs.append(str(pfad))
    if not dirs:
        return {}
    cache = workdir / "fccache"
    cache.mkdir(parents=True, exist_ok=True)
    conf = workdir / "fonts.conf"
    zeilen = ['<?xml version="1.0"?>', '<!DOCTYPE fontconfig SYSTEM "fonts.dtd">', "<fontconfig>"]
    zeilen += [f"  <dir>{d}</dir>" for d in dirs]
    zeilen += [f"  <cachedir>{cache}</cachedir>", "</fontconfig>", ""]
    conf.write_text("\n".join(zeilen), encoding="utf-8")
    return {"FONTCONFIG_FILE": str(conf)}


def _soffice_env(extra_font_dirs: list[Path] | None, workdir: Path) -> dict:
    """The full environment for one headless run: fonts plus a profile of its own.

    The separate profile is not cosmetic. A headless run against the user's own LibreOffice profile
    refuses to start while they have LibreOffice open, and it also writes into settings they own.
    """
    import os
    env = dict(os.environ)
    env.update(_fontconfig_for(extra_font_dirs, workdir))
    profil = (workdir / "loprofile").resolve()
    profil.mkdir(parents=True, exist_ok=True)
    # Resolved, because a path through a symlink is a different string for the same place and this
    # one is handed to another process as a URL.
    env["_ZANMAI_LO_PROFILE"] = f"-env:UserInstallation=file://{profil}"
    return env


def render(deck_path: Path, into: Path, dpi: int = 110, font_dirs: list[Path] | None = None) -> int:
    """A picture of every slide, headless, on any platform.

    This exists because the alternative was nothing. `qlmanage` renders the first slide only,
    on macOS only, and puts itself in the dock while it runs; outside macOS there was no way
    to look at a deck at all, and a check that cannot look is a check that trusts its own
    description. Measured while building the wireframe library: 58 slides rendered
    in about 6 seconds, and the render found ten faults that every other check here reported
    clean -- shapes overlapping, a marker breaking over two lines, banding that was white on
    white. None of those wrap, none of them is a broken reference, and none of them shows up
    in the XML as anything unusual.

    **What it is and is not.** LibreOffice lays a deck out again rather than reproducing
    PowerPoint exactly, so this answers "does the arrangement work, does the text sit in its
    box, is anything on top of anything else" and not "is this pixel-identical to what the
    customer sees". Known gaps, verified: a shape inherited from the layout and an `outerShdw`
    may render differently than PowerPoint does. For the questions above that is enough, and
    for the questions below it it was never the tool.
    """
    exe = _soffice()
    if exe is None:
        print(f"slide-library: LibreOffice is not on this machine, so a deck cannot be looked "
              f"at here. To make it possible: {RENDER_HINT}", file=sys.stderr)
        return 1
    import subprocess
    import tempfile
    into.mkdir(parents=True, exist_ok=True)
    filter_arg = ('pdf:impress_pdf_Export:{"ExportHiddenSlides":{"type":"boolean",'
                  '"value":"true"}}')
    with tempfile.TemporaryDirectory() as tmp:
        # Hidden slides are skipped by the export by default and nothing says so: measured on a
        # real deck of 27 slides where 16 were hidden and 11 came out, with no warning at all.
        env = _soffice_env(font_dirs, Path(tmp))
        lauf = subprocess.run([exe, "--headless", env["_ZANMAI_LO_PROFILE"],
                               "--convert-to", filter_arg, "--outdir", tmp,
                               str(deck_path)], capture_output=True, text=True, env=env)
        pdfs = list(Path(tmp).glob("*.pdf"))
        if not pdfs:
            print(f"slide-library: LibreOffice produced no PDF for {deck_path.name}. "
                  f"{lauf.stderr.strip() or lauf.stdout.strip()}", file=sys.stderr)
            return 1
        stamm = into / deck_path.stem
        try:
            import fitz  # noqa: F401 -- PyMuPDF, when it happens to be there
        except ImportError:
            fitz = None
        import shutil as _shutil
        if _shutil.which("pdftoppm"):
            subprocess.run(["pdftoppm", "-png", "-r", str(dpi), str(pdfs[0]), str(stamm)],
                           check=True)
        elif fitz is not None:
            doc = fitz.open(str(pdfs[0]))
            for i, seite in enumerate(doc, start=1):
                seite.get_pixmap(dpi=dpi).save(f"{stamm}-{i:02d}.png")
        else:
            ziel = into / (deck_path.stem + ".pdf")
            _shutil.copy(str(pdfs[0]), str(ziel))
            print(f"rendered {deck_path.name} -> {ziel} (PDF only: neither pdftoppm nor PyMuPDF "
                  f"is here to cut it into pictures)")
            return 0
    bilder = sorted(into.glob(f"{deck_path.stem}-*.png"))
    print(f"rendered {len(bilder)} slide(s) of {deck_path.name} into {into}")
    for bild in bilder:
        print(f"  {bild}")
    return 0


def brand_master(deck, scheme_name: str | None = None):
    """The master to build on, chosen by the colour scheme its own theme carries.

    A deck can hold several masters and only one of them be the brand: measured on a
    real company template with three, two of which were stock Office. Picking a layout by name
    landed on the Office master and produced a slide that looked finished in the wrong colours,
    with nothing to show it was wrong. Without a name, the first master whose scheme is not
    called "Office" wins, and failing that the first one.
    """
    kandidaten = []
    for master in deck.slide_masters:
        name = None
        for rel in master.part.rels.values():
            if "theme" in rel.reltype:
                xml = rel.target_part.blob.decode("utf-8", "ignore")
                treffer = re.search(r'<a:clrScheme name="([^"]*)"', xml)
                name = treffer.group(1) if treffer else None
                break
        kandidaten.append((master, name or ""))
        if scheme_name and scheme_name.lower() in (name or "").lower():
            return master, name
    if scheme_name:
        raise SystemExit(f"slide-library: no master carries a theme named like {scheme_name!r}; "
                         f"this deck has {[n for _m, n in kandidaten]}")
    for master, name in kandidaten:
        if name and name.lower() != "office":
            return master, name
    return kandidaten[0] if kandidaten else (None, None)


def migrate(source_path: Path, slide_no: int, target_path: Path, out: Path,
            layout_name: str | None = None, scheme: str | None = None,
            title: str | None = None, replace: bool = False, keep_shapes: bool = False,
            brand_from: Path | None = None) -> int:
    """Put one slide into another deck, so it belongs there instead of bringing its own world.

    `clone` copies inside one deck: same master, same theme, and a colour difference would be a
    fault. This is the opposite case. The slide arrives in a file that already carries a brand,
    and what comes across is the arrangement: shapes, geometry and text. Master, layouts, theme,
    fonts and whatever the layout paints -- a logo, a background shape -- come from the target,
    because the slide is added to a copy of the target rather than the theme being copied into
    the source.

    Two things this reports rather than silently fixing. A shape that names a colour outright
    instead of a theme role keeps that colour, and is listed: only a person can say which brand
    role a literal `#2E86AB` was standing in for. And a title placeholder in the chosen layout is
    filled if `title` is given, because a title written into a placeholder inherits the brand's
    own type, while one drawn as a text box does not.
    """
    warnungen: list[str] = []
    quelle = Presentation(str(source_path))
    if not 1 <= slide_no <= len(quelle.slides._sldIdLst):
        print(f"slide-library: {source_path.name} has no slide {slide_no}", file=sys.stderr)
        return 1
    folie = quelle.slides[slide_no - 1]

    # Where the brand is measured. The target is the default and the wrong default in the one case
    # that matters most: a deck built set by set starts empty, so the brand lives in another file
    # and everything measured here finds nothing. Seen in practice, where the
    # first sets came out square and the later ones round, with nothing said either way.
    marke = brand_from or target_path
    marken_warnung = _theme_carries_the_brand(marke)
    if marken_warnung:
        warnungen.append(marken_warnung)
    ziel = Presentation(str(target_path))
    lst = ziel.slides._sldIdLst
    vorhanden = len(lst)
    # This used to empty the target unconditionally, and it cost real work: two migrates in a row
    # do not accumulate that way, the second one throws the first one's slide away. Reported from
    # the field, with the built file lost because `--out` pointed at it. `--into` reads
    # as putting something in, and composing a deck out of several patterns is the actual use, so
    # appending is now what happens and emptying is a flag.
    if replace:
        for sld in list(lst):
            ziel.part.drop_rel(sld.rId)
            lst.remove(sld)
    master, scheme_name = brand_master(ziel, scheme)
    if master is None:
        print(f"slide-library: {target_path.name} carries no master", file=sys.stderr)
        return 1
    layouts = list(master.slide_layouts)
    if layout_name:
        passend = [L for L in layouts if L.name.lower() == layout_name.lower()]
        if not passend:
            print(f"slide-library: {target_path.name} has no layout {layout_name!r} on the "
                  f"{scheme_name!r} master; it has {[L.name for L in layouts]}", file=sys.stderr)
            return 1
        layout = passend[0]
    else:
        layout = next((L for L in layouts if L.name.lower() in ("leer", "blank", "inhalt 1")),
                      layouts[0])
    neu = ziel.slides.add_slide(layout)

    titel_ph = None
    for ph in list(neu.placeholders):
        if ph.placeholder_format.idx == 0 and title:
            titel_ph = ph
        else:
            ph._element.getparent().remove(ph._element)
    if titel_ph is not None:
        titel_ph.text_frame.text = title

    # A slide that already carries its own title text box, plus a title written into the
    # layout's placeholder, paints both on top of each other. Reported, not silently dropped:
    # which of the two is the real title is not this command's call.
    if titel_ph is not None:
        oben = [sh for sh in folie.shapes
                if sh.has_text_frame and text_of(sh)
                and inches(sh.top) < inches(titel_ph.top) + inches(titel_ph.height)]
        if oben:
            warnungen.append(
                "the source slide already carries text where the layout's title sits (%s); "
                "both will paint. Drop --title, or remove that shape from the source."
                % ", ".join(repr(text_of(sh)[:20]) for sh in oben[:2]))

    fest = []
    for shape in folie.shapes:
        kopie = copy.deepcopy(shape._element)
        _carry_relationships(neu, folie, kopie)
        neu.shapes._spTree.append(kopie)
        for treffer in re.findall(r'srgbClr val="([0-9A-Fa-f]{6})"',
                                  kopie.xml if hasattr(kopie, "xml") else ""):
            fest.append((shape.name, treffer))

    # Two things a theme swap changes that no colour check sees.
    quell_schrift = _theme_fonts(source_path).get("minor")
    ziel_schrift = None
    for rel in master.part.rels.values():
        if "theme" in rel.reltype:
            xml = rel.target_part.blob.decode("utf-8", "ignore")
            treffer = re.search(r"<a:minorFont><a:latin typeface=\"([^\"]*)\"", xml)
            ziel_schrift = treffer.group(1) if treffer else None
            break
    if quell_schrift and ziel_schrift and quell_schrift != ziel_schrift:
        vor, nach = face_width_factor(quell_schrift), face_width_factor(ziel_schrift)
        if vor and nach and abs(nach - vor) > 0.02:
            warnungen.append(
                "the body face changes from %s to %s, which runs %+.0f%% wider, so every slot "
                "holds about %.0f%% less: run overflow-check on the result and shorten what no "
                "longer fits." % (quell_schrift, ziel_schrift, (nach / vor - 1) * 100,
                                  (1 - vor / nach) * 100))
    # A brand may put the same value on two roles. Where a slide told two things apart by using
    # both, that distinction is gone and nothing in the file says so.
    rollen = {}
    for rel in master.part.rels.values():
        if "theme" in rel.reltype:
            xml = rel.target_part.blob.decode("utf-8", "ignore")
            schema = re.search(r"<a:clrScheme.*?</a:clrScheme>", xml, re.S)
            if schema:
                for rolle, srgb, sys_ in re.findall(
                        r"<a:(lt1|lt2|dk1|dk2|accent[1-6])><a:(?:srgbClr val=\"([0-9A-Fa-f]{6})\""
                        r"|sysClr val=\"[^\"]*\" lastClr=\"([0-9A-Fa-f]{6})\")", schema.group(0)):
                    rollen.setdefault((srgb or sys_).upper(), []).append(rolle)
            break
    doppelt = {wert: r for wert, r in rollen.items() if len(r) > 1}
    for wert, r in doppelt.items():
        warnungen.append(
            "this brand puts the same colour #%s on %s, so anything that told those two apart "
            "now reads as one." % (wert, " and ".join(r)))

    beispiele = {} if keep_shapes else _brand_examples(marke)
    uebernommen, look_detail = _apply_brand_look(neu, beispiele) if beispiele else (0, [])
    if not keep_shapes and not beispiele:
        warnungen.append(
            f"no brand example found in {marke.name}: it carries too few shapes of a kind to tell "
            f"what is usual, so the migrated slide keeps the wireframe's own look. Point "
            f"`--brand-from` at a deck that already carries the brand.")

    if not guarded_save(ziel, out):
        return 1
    print(f"migrated slide {slide_no} of {source_path.name} into {target_path.name} "
          f"-> {out} (master {scheme_name!r}, layout {layout.name!r})")
    if uebernommen:
        quellen = ", ".join(f"{art} from {n} of {gesamt} in {marke.name}"
                            for art, (_el, n, gesamt) in sorted(beispiele.items()))
        print(f"  gave {uebernommen} shape(s) the brand's own look ({', '.join(look_detail)}), taken "
              f"whole from a real shape of each kind rather than property by property: {quellen}. "
              f"A theme carries colour and type and nothing about form, fill, line or shadow, so "
              f"without this the wireframe's own look comes across. `--keep-shapes` leaves it.")
    moebel = _wireframe_furniture(neu)
    if moebel:
        warnungen.append(
            "the wireframe brought %d textless bar(s) or marker(s) across: %s. A wireframe draws "
            "its own aids to show where things go, and they are the sketch's language rather than "
            "the piece's. Keep one only where the brand itself uses that element in that role."
            % (len(moebel), ", ".join(moebel[:4]) + (" ..." if len(moebel) > 4 else "")))
    if replace and vorhanden:
        print(f"  --replace: the {vorhanden} slide(s) the target carried were dropped, "
              f"{out.name} holds the migrated one alone")
    else:
        print(f"  appended: {out.name} now holds {len(lst)} slide(s), the migrated one last")
    if fest:
        gezaehlt = {}
        for name, wert in fest:
            gezaehlt.setdefault(wert, []).append(name)
        print(f"  {len(gezaehlt)} literal colour(s) came across unchanged; only a person can say "
              f"which brand role each stood for:")
        for wert, namen in sorted(gezaehlt.items(), key=lambda kv: -len(kv[1])):
            print(f"    #{wert}  {len(namen)}x  e.g. {namen[0]}")
    else:
        print("  no literal colours: every fill and every run named a theme role, so the target's "
              "theme colours all of it")
    for warnung in warnungen:
        print(f"  CHECK: {warnung}")
    return 0


_R_ID = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
_SLIDE_REL = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide"


def suggest(library: Path, elements: int | None, order: str | None, movement: str | None) -> int:
    """Which patterns carry a content of this shape. Narrows 57 down to a handful.

    Not a decision, a shortlist. The choice between the candidates is a judgement about this
    content, and it stays with the person making it; what this removes is the part that was never
    judgement, reading a few of 57 prose descriptions and taking the first that sounds plausible.
    Seen in practice: the pattern used for one piece was reused for a piece whose
    content had a different shape, and a plausible reason was written for both.

    The shape of the content is described first and on its own: how many things there are, whether
    they are of equal rank or in an order, and whether something moves from A to B. Then the
    patterns are matched against that, rather than the content being read out of a pattern.
    """
    datei = library / "library.json"
    if not datei.is_file():
        print(f"slide-library: no library.json in {library}", file=sys.stderr)
        return 2
    daten = json.loads(datei.read_text(encoding="utf-8"))
    muster = daten.get("patterns") or daten.get("slides") or []
    treffer = []
    for eintrag in muster:
        if "order" not in eintrag:
            continue
        punkte, warum = 0, []
        if order:
            if eintrag["order"] != order:
                continue
            punkte += 2
            warum.append(order)
        if movement is not None:
            will = movement.lower() in ("yes", "y", "true", "1", "ja")
            if bool(eintrag.get("movement")) != will:
                continue
            punkte += 2
            warum.append("movement" if will else "no movement")
        if elements is not None:
            von, bis = (eintrag.get("elements") or [1, 1])[:2]
            if not von <= elements <= bis:
                continue
            punkte += 1
            warum.append(f"{elements} fits {von}-{bis}")
        treffer.append((punkte, eintrag, warum))
    if not treffer:
        print("no pattern carries that shape. Either the content has a shape the library does not "
              "hold, which is what composing is for, or one of the three answers is too narrow.")
        return 1
    treffer.sort(key=lambda x: (-x[0], x[1]["id"]))
    print(f"{len(treffer)} pattern(s) carry that shape:\n")
    for _punkte, eintrag, warum in treffer[:8]:
        von, bis = (eintrag.get("elements") or [1, 1])[:2]
        print(f"  {eintrag['id']:<24} slide {eintrag.get('slide'):<3} {von}-{bis} elements, "
              f"{eintrag['order']}{', movement' if eintrag.get('movement') else ''}")
        print(f"    {eintrag.get('structure', '')[:100]}")
        print(f"    carrier: {eintrag.get('carrier', '?')}")
    if len(treffer) > 8:
        print(f"\n  ... and {len(treffer) - 8} more")
    print("\nNow judge between these against the content. The shortlist is not the decision.")
    return 0


def keep(deck_path: Path, slide_no: int, brand_dir: Path, slug: str, source: str | None) -> int:
    """Put an approved slide into the brand's own library, so the next one of its kind is a copy.

    This is the step that was missing, and its absence is measurable. A first build of two slides
    took 26 minutes and 76 tool calls; the correction of the same two took 2.5 minutes and 11. The
    difference was not skill or tooling, it was that the way was already known the second time. With
    nothing kept, every run pays the first price again: reading the bundle, measuring the brand,
    deriving the route, and then writing a one-off script anyway.

    What is kept is the slide itself, not a description of it. The slide carries every value someone
    would otherwise have to write down: corner, fill, line, shadow, spacing, type. A note beside it
    holds only what the file cannot say, namely where it came from and what its fillable places are.

    The trigger is approval, never the build. What the user rejected has no business in the brand.
    """
    ziel_ordner = brand_dir / "slides"
    ziel_ordner.mkdir(parents=True, exist_ok=True)
    ziel = ziel_ordner / f"{slug}.pptx"
    if extract(deck_path, [slide_no], ziel) != 0:
        return 1

    deck = Presentation(str(ziel))
    slots = slot_names(deck.slides[0])
    stand = datetime.now(timezone.utc).strftime("%Y-%m-%d")
    zeilen = [f"# {slug}", "",
              f"Approved slide, kept {stand}. Take a copy of `{slug}.pptx` and fill it; do not "
              f"rebuild this shape.", ""]
    if source:
        zeilen += [f"Came from: {source}", ""]
    zeilen += ["| slot | holds | now |", "|---|---|---|"]
    for name, slot in slots.items():
        zeilen.append(f"| `{name}` | {slot['capacity']} chars | {slot['text'][:40]!r} |")
    zeilen += ["", "```", f"slide-library.py extract {slug}.pptx --slides 1 --out <new.pptx>",
               "slide-library.py fill <new.pptx> --texts <texts.json> --out <done.pptx>", "```", ""]
    (ziel_ordner / f"{slug}.md").write_text("\n".join(zeilen), encoding="utf-8")

    eintraege = sorted(f.stem for f in ziel_ordner.glob("*.pptx"))
    index = ["# Brand slides", "",
             "Slides the user approved. **Look here before building anything**: taking one of these",
             "and swapping its text is the cheapest route there is, and the only one where the look",
             "cannot drift, because nothing is redrawn.", ""]
    for eintrag in eintraege:
        note = ziel_ordner / f"{eintrag}.md"
        erste = ""
        if note.is_file():
            for zeile in note.read_text(encoding="utf-8").splitlines():
                if zeile.startswith("Approved slide"):
                    erste = zeile
                    break
        index.append(f"- [`{eintrag}`]({eintrag}.md): {erste}")
    (ziel_ordner / "INDEX.md").write_text("\n".join(index) + "\n", encoding="utf-8")

    print(f"kept slide {slide_no} as {ziel} ({len(slots)} fillable place(s))")
    print(f"  {len(eintraege)} slide(s) in this brand's library now. Where a later piece needs this "
          f"same shape of content, it is a copy of this file plus `fill` rather than a build. Where "
          f"it needs a different shape, this file is not the answer and its cheapness is not a "
          f"reason to make it one.")
    return 0


def extract(deck_path: Path, slides: list[int], out: Path) -> int:
    """Lift whole slides out of a deck, keeping their own XML untouched.

    This is the fastest form of Match written as one command: the slides that carry the shape of the
    content stay exactly as they are, geometry byte-identical, and everything else goes. Doing it by
    hand is three steps, and the third one is easy to miss.

    That third step is what this exists for. Deleting a navigation button's shape is not enough: the
    slide's own relationship to the slide it pointed at stays behind, and it keeps every part that
    slide reaches alive in the file. Measured: 7.1 MB against 4.1 MB after
    the relationship was cut, fifteen foreign slide parts held by one leftover link. PowerPoint also
    calls a file damaged over a reference into nothing, so this is not only about size.
    """
    deck = Presentation(str(deck_path))
    gesamt = len(deck.slides._sldIdLst)
    ungueltig = [nr for nr in slides if not 1 <= nr <= gesamt]
    if ungueltig:
        print(f"slide-library: {deck_path.name} has {gesamt} slide(s), asked for "
              f"{', '.join(str(nr) for nr in ungueltig)}", file=sys.stderr)
        return 2
    if not slides:
        print("slide-library: name at least one slide", file=sys.stderr)
        return 2

    lst = deck.slides._sldIdLst
    alle = list(lst)
    behalten = [alle[nr - 1] for nr in slides]        # in the order asked for
    # Re-order first, then drop: a slide moved out of the list keeps its part, so the order the
    # caller asked for is the order the new deck has.
    for sld in behalten:
        lst.remove(sld)
        lst.append(sld)
    for sld in alle:
        if sld not in behalten:
            deck.part.drop_rel(sld.get(_R_ID))
            lst.remove(sld)

    # The links the kept slides carry into slides that are gone now. Both halves have to go: the
    # relationship, which is what holds the foreign slide's parts alive, and the `hlinkClick` in the
    # shape that names it, because a shape pointing at a relationship that is no longer there is
    # exactly the reference into nothing PowerPoint calls damage. The shape itself stays; whether a
    # navigation button belongs on the slide is not a decision this command may take.
    gekappt = 0
    behalten_teile = [s.part for s in deck.slides]
    for slide in deck.slides:
        for rel_id, rel in list(slide.part.rels.items()):
            if rel.reltype != _SLIDE_REL:
                continue
            ziel_teil = rel.target_part if not rel.is_external else None
            if ziel_teil is not None and ziel_teil in behalten_teile:
                continue
            for hlink in slide.shapes._spTree.iter(qn("a:hlinkClick")):
                if hlink.get(_R_ID) == rel_id:
                    hlink.getparent().remove(hlink)
            slide.part.drop_rel(rel_id)
            gekappt += 1

    if not guarded_save(deck, out):
        return 1
    vorher = deck_path.stat().st_size / 1048576
    nachher = out.stat().st_size / 1048576
    print(f"extracted slide(s) {', '.join(str(nr) for nr in slides)} of {deck_path.name} -> {out} "
          f"({nachher:.1f} MB, was {vorher:.1f} MB)")
    if gekappt:
        print(f"  cut {gekappt} link(s) into slides that are not in the new file; without that they "
              f"would have held those slides' parts alive and PowerPoint would call the file damaged")
    lose = dangling_refs(out)
    for eintrag in lose:
        print(f"FAIL: {eintrag}", file=sys.stderr)
    return 1 if lose else 0


def _visible_count(folie) -> dict:
    """What a slide is made of: shapes, filled areas, text places, pictures."""
    formen = gefuellt = texte = bilder = 0
    for shape, (_l, _t, breite, hoehe) in _every_leaf(folie.shapes):
        if not breite or not hoehe:
            continue
        formen += 1
        if shape._element.tag == qn("p:pic"):
            bilder += 1
            continue
        if shape.has_text_frame and text_of(shape):
            texte += 1
        spPr = shape._element.find(qn("p:spPr"))
        if spPr is not None and (spPr.find(qn("a:solidFill")) is not None
                                 or spPr.find(qn("a:gradFill")) is not None):
            gefuellt += 1
    return {"shapes": formen, "filled": gefuellt, "text": texte, "pictures": bilder}


def structure_check(deck_path: Path, slide_no: int, gegen: str, toleranz: float = 0.0,
                    intended: str | None = None) -> int:
    """Is the built slide still carrying the structure of the pattern it was built from.

    Every other check here asks whether what is there is right. None of them asks whether something
    that should be there is missing, and that is the failure they let through: a hub with four
    satellites and no hub, a stage band with the stages gone, two text columns with the areas that
    separated them dropped. All five checks reported clean on all three, correctly, because nothing
    that remained was wrong. Seen in practice, found by a person on the sixteenth
    slide, after fifteen had been reported green.

    **A difference is not a fault.** A wireframe is a starting point, not a template to be copied:
    dropping a part of it, adding one, rebuilding the arrangement around the content is what it is
    there for, and a check that forbade that would turn the sketch into a specification. What this
    catches is the other case, the one nobody chose: a shape that fell out during the build, whose
    absence leaves no trace anywhere else. The two look identical from here, so the difference is
    reported and the person building says which it was, with `--intended`.

    **The default tolerance is zero, and that is deliberate.** The first version allowed a quarter,
    and both slides this was built for came through it clean: a hub missing its centre still holds
    six of eight filled areas, and the one that is gone is the thing all four lines point at. A
    proportion says nothing about whether what is missing was load-bearing.
    """
    if ":" not in gegen:
        print(f"slide-library: --against takes <deck>:<slide>, got {gegen!r}", file=sys.stderr)
        return 2
    quell_pfad, _, quell_nr = gegen.rpartition(":")
    quelle = Path(quell_pfad)
    if not quelle.is_file():
        print(f"slide-library: no deck at {quelle}", file=sys.stderr)
        return 2
    try:
        quell_index = int(quell_nr)
    except ValueError:
        print(f"slide-library: --against needs a slide number, got {quell_nr!r}", file=sys.stderr)
        return 2

    gebaut = Presentation(str(deck_path))
    muster = Presentation(str(quelle))
    if not 1 <= slide_no <= len(gebaut.slides._sldIdLst):
        print(f"slide-library: {deck_path.name} has no slide {slide_no}", file=sys.stderr)
        return 2
    if not 1 <= quell_index <= len(muster.slides._sldIdLst):
        print(f"slide-library: {quelle.name} has no slide {quell_index}", file=sys.stderr)
        return 2

    ist = _visible_count(gebaut.slides[slide_no - 1])
    soll = _visible_count(muster.slides[quell_index - 1])
    fehlend = []
    for art in ("shapes", "filled", "text", "pictures"):
        if soll[art] and ist[art] < soll[art] * (1 - toleranz):
            fehlend.append(f"{art}: {ist[art]} against {soll[art]} in the pattern")
    # What the pattern is, in its own words, so a report can be judged instead of just counted.
    beschreibung, traeger = "", ""
    bibliothek = quelle.parent / "library.json"
    if bibliothek.is_file():
        try:
            eintraege = json.loads(bibliothek.read_text(encoding="utf-8"))
            for eintrag in (eintraege.get("patterns") or eintraege.get("slides") or []):
                if eintrag.get("slide") == quell_index:
                    beschreibung = eintrag.get("structure", "")
                    traeger = eintrag.get("carrier", "")
                    break
        except (ValueError, OSError):
            beschreibung = ""

    print(f"slide {slide_no} of {deck_path.name} against slide {quell_index} of {quelle.name}:")
    if beschreibung:
        print(f"  pattern: {beschreibung}")
    if traeger:
        print(f"  carrier: {traeger}")
    for art in ("shapes", "filled", "text", "pictures"):
        print(f"  {art:<9} {ist[art]:>4}   pattern {soll[art]:>4}")
    # The carrier is the shape the pattern is about. Its loss is never a matter of proportion, and
    # `--intended` does not cover it either: without the centre a hub is four labels, and that is a
    # different pattern, not this one changed. Named so the person looks at exactly that.
    if fehlend and traeger:
        print(f"NOTE: the part this pattern is built around is {traeger}. If that is what is "
              f"missing, the slide no longer carries the pattern, whatever the counts say.")
    if fehlend and intended:
        for eintrag in fehlend:
            print(f"  differs: {eintrag}")
        print(f"stated as intended: {intended}")
        print("A wireframe is a starting point, so a deliberate difference is the normal case. "
              "Recorded, not failed.")
        return 0
    if fehlend:
        for eintrag in fehlend:
            print(f"FAIL: {eintrag}")
        print(f"{len(fehlend)} difference(s) from the pattern. Look at the slide and ask which of "
              f"the two this is: a shape that fell out during the build, which leaves no trace any "
              f"other check can see, or a part you deliberately did without. A wireframe is a "
              f"starting point and may be changed freely; say so with --intended \"<why>\" and this "
              f"passes. What must not happen is the difference going unnoticed.")
        return 1
    print("the pattern arrived whole")
    return 0


def leftover_check(deck_path: Path) -> int:
    """What the file carries that is not the content: comments, speaker notes, animations, authors.

    A cloned slide brings its whole history along, and none of it shows in a render or in any
    geometry check. A comment written during review can stand on a finished slide, next to a notes
    page and the source's animations, and be found only by opening the file. This is the kind of defect that becomes visible at the recipient.

    Reported, never removed. Speaker notes are often wanted, an animation may be the point, and
    which of them belongs in a handover is not a decision a check may take.
    """
    import zipfile
    befunde: list[str] = []
    try:
        with zipfile.ZipFile(deck_path) as z:
            namen = z.namelist()

            kommentare = [nm for nm in namen if re.match(r"ppt/(comments|threadedComments)/", nm)]
            if kommentare:
                autoren = set()
                for nm in [x for x in namen if x.endswith(("commentAuthors.xml", "authors.xml"))] + kommentare:
                    autoren |= set(re.findall(r'name="([^"]{2,})"',
                                              z.read(nm).decode("utf-8", "ignore")))
                # A threaded comment names its author by GUID and the name lives in `authors.xml`;
                # printing the GUID would be a finding nobody can act on.
                autoren = {a for a in autoren if not re.fullmatch(r"\{?[0-9A-Fa-f-]{30,}\}?", a)}
                wer = f", by {', '.join(sorted(autoren)[:3])}" if autoren else ""
                befunde.append(f"{len(kommentare)} comment part(s) in the file{wer}")

            notizen = []
            for nm in sorted(x for x in namen if re.match(r"ppt/notesSlides/notesSlide\d+\.xml$", x)):
                text = " ".join(re.findall(r"<a:t>([^<]*)</a:t>",
                                           z.read(nm).decode("utf-8", "ignore"))).strip()
                # A notes page carrying only its own slide number is what PowerPoint writes by
                # itself; calling that a leftover would make the check noise.
                if text and not text.isdigit():
                    notizen.append((nm.rsplit("/", 1)[-1], text[:60]))
            if notizen:
                befunde.append(f"{len(notizen)} speaker-notes page(s) with text, e.g. "
                               f"{notizen[0][0]}: {notizen[0][1]!r}")

            animiert = []
            for nm in sorted(x for x in namen if re.match(r"ppt/slides/slide\d+\.xml$", x)):
                xml = z.read(nm).decode("utf-8", "ignore")
                if re.search(r"<p:(par|seq|anim|animEffect|animMotion)\b", xml):
                    animiert.append(nm.rsplit("/", 1)[-1])
            if animiert:
                befunde.append(f"{len(animiert)} slide(s) carry animation: {', '.join(animiert[:4])}"
                               + (" …" if len(animiert) > 4 else ""))

            if "docProps/core.xml" in namen:
                kern = z.read("docProps/core.xml").decode("utf-8", "ignore")
                spuren = [f"{feld}={wert!r}" for feld, wert in
                          re.findall(r"<(?:dc|cp):(creator|lastModifiedBy)>([^<]+)<", kern) if wert]
                if spuren:
                    befunde.append("the file's properties name " + ", ".join(spuren))
    except (OSError, KeyError, zipfile.BadZipFile) as fehler:
        print(f"slide-library: cannot read {deck_path}: {fehler}", file=sys.stderr)
        return 2

    for befund in befunde:
        print(f"FAIL: {befund}")
    print(f"{len(befunde)} leftover(s) in {deck_path.name}. Nothing was removed; which of these "
          f"belongs in a handover is a person's call.")
    return 1 if befunde else 0


def dangling_refs(deck_path: Path) -> list[str]:
    """Every relationship id a slide points at that its own rels cannot resolve.

    Read from the written file rather than from the object in memory, because that is the file the
    user opens. A reference into nothing is always a defect, whatever produced it.
    """
    import zipfile
    problems = []
    with zipfile.ZipFile(deck_path) as z:
        namen = z.namelist()
        for slide in sorted(n for n in namen if re.match(r"ppt/slides/slide\d+\.xml$", n)):
            xml = z.read(slide).decode("utf-8", "ignore")
            ids = set(re.findall(r'r:(?:embed|link|id)="([^"]+)"', xml))
            rels_name = slide.replace("slides/", "slides/_rels/") + ".rels"
            rels = z.read(rels_name).decode("utf-8", "ignore") if rels_name in namen else ""
            vorhanden = set(re.findall(r'Id="([^"]+)"', rels))
            for fehlend in sorted(ids - vorhanden):
                problems.append(f"{Path(slide).name} points at {fehlend}, which its rels do not carry")
    return problems


# What the first bytes of a picture look like, per format. A part that is present, correctly
# referenced and empty passes every other check here: the reference resolves, the XML is valid, and
# a render looks right because the renderer quietly substitutes. PowerPoint is the only thing that
# notices, and it notices by offering to repair the file. That happened on a finished deck, and the
# only way anyone found out was by opening it, which is exactly the thing that must not be the way.
_MAGIC = {
    "png": (b"\x89PNG\r\n\x1a\n",),
    "jpg": (b"\xff\xd8\xff",),
    "jpeg": (b"\xff\xd8\xff",),
    "gif": (b"GIF87a", b"GIF89a"),
    "bmp": (b"BM",),
    "tif": (b"II*\x00", b"MM\x00*"),
    "tiff": (b"II*\x00", b"MM\x00*"),
    "emf": (b"\x01\x00\x00\x00",),
    "wmf": (b"\xd7\xcd\xc6\x9a", b"\x01\x00\x09\x00"),
    "svg": (b"<svg", b"<?xml"),
    "webp": (b"RIFF",),
}


def media_check(deck_path: Path, verbose: bool = False) -> int:
    """Every picture a deck points at, checked for being an actual file.

    Three questions, and only the first two were ever asked. `refs-check` answers "does the
    reference point somewhere". The schema check answers "is the order of elements legal". Nobody
    answered "is the thing pointed at a real picture", and that is the one that reached a user: a
    media part correctly referenced and zero bytes long, valid XML, a render that looked right, and
    PowerPoint offering to repair the file on open.

    Reading the first bytes rather than trusting the extension, because that is the difference
    between a file named `.png` and a PNG.
    """
    import zipfile

    problems: list[str] = []
    geprueft = 0
    with zipfile.ZipFile(deck_path) as z:
        namen = set(z.namelist())
        groessen = {i.filename: i.file_size for i in z.infolist()}
        rels = [n for n in namen if re.match(r"ppt/(slides|slideLayouts|slideMasters)/_rels/.*\.rels$", n)]
        gesehen: set[str] = set()
        for rel in sorted(rels):
            xml = z.read(rel).decode("utf-8", "ignore")
            for rid, ziel in re.findall(r'Id="([^"]+)"[^>]*Type="[^"]*/image"[^>]*Target="([^"]+)"', xml):
                # A target is relative to the part's own folder, `../media/image1.png` from a slide.
                basis = Path(rel).parent.parent
                teil = str((basis / ziel).as_posix()).replace("/./", "/")
                while "/../" in teil:
                    teil = re.sub(r"[^/]+/\.\./", "", teil, count=1)
                quelle = Path(rel).name.replace(".rels", "")
                if teil in gesehen:
                    continue
                gesehen.add(teil)
                geprueft += 1
                if teil not in namen:
                    problems.append(f"{quelle} points at {teil}, which is not in the file at all")
                    continue
                groesse = groessen.get(teil, 0)
                if groesse == 0:
                    problems.append(f"{teil} is in the file and is empty (0 bytes); "
                                    f"a render substitutes and PowerPoint offers to repair")
                    continue
                endung = Path(teil).suffix.lstrip(".").lower()
                erwartet = _MAGIC.get(endung)
                if erwartet:
                    kopf = z.read(teil)[:16]
                    if not any(kopf.startswith(m) for m in erwartet):
                        problems.append(f"{teil} calls itself {endung} and does not start like one "
                                        f"({kopf[:4]!r}, {groesse} bytes)")
                        continue
                if verbose:
                    print(f"  ok {teil}  {groesse} bytes  {endung or 'no extension'}")
    for problem in problems:
        print(f"FAIL: {problem}")
    print(f"{geprueft} picture(s) checked, {len(problems)} unusable"
          + _nichts_geprueft(geprueft))
    return 1 if problems else 0


_TABLE_SLOT = re.compile(r"^table(\d+)\.r(\d+)c(\d+)$")


def _locate_slot_frame(slide, name: str, template_text: str, vergeben: set | None = None):
    """The text frame a slot name resolves to in this (cloned) slide: a table cell, addressed by
    position since cell text need not be unique, or a shape, addressed by matching its current text
    against what harvest recorded, exactly as before tables were slots too.

    Groups are walked, because `slot_names` names shapes inside them. `vergeben` holds the shapes
    an earlier slot in the same run already claimed: two rubrics of a template often carry the
    same placeholder wording, and without this the second write would land on the first one again
    and leave the second standing untouched.
    """
    table_match = _TABLE_SLOT.match(name)
    if table_match:
        wanted_table, row, col = (int(g) for g in table_match.groups())
        seen = 0
        for shape, _rect in _every_leaf(slide.shapes):
            if not getattr(shape, "has_table", False):
                continue
            seen += 1
            if seen == wanted_table:
                return shape.table.cell(row - 1, col - 1).text_frame
        return None
    for shape, _rect in _every_leaf(slide.shapes):
        if shape.has_text_frame and text_of(shape) == template_text:
            if vergeben is not None:
                if id(shape) in vergeben:
                    continue
                vergeben.add(id(shape))
            return shape.text_frame
    return None


def _fill_text_frame(frame, lines: list) -> str | None:
    """Replace a text frame's content with `lines`, one paragraph per line, keeping the template's
    own run formatting. Returns an error string, or None on success."""
    paragraphs = frame.paragraphs
    template_run = None
    for paragraph in paragraphs:
        if paragraph.runs:
            template_run = paragraph.runs[0]
            break
    if template_run is None:
        return "carries no styled run to follow"
    keep = copy.deepcopy(template_run._r)
    keep_paragraph = copy.deepcopy(paragraphs[0]._p)
    frame.clear()
    for index, line in enumerate(lines):
        if index == 0:
            paragraph = frame.paragraphs[0]
        else:
            new_p = copy.deepcopy(keep_paragraph)
            for run in new_p.findall(
                    "{http://schemas.openxmlformats.org/drawingml/2006/main}r"):
                new_p.remove(run)
            frame._txBody.append(new_p)
            paragraph = frame.paragraphs[-1]
        # One paragraph can carry several runs with different weight. A rubric like an objection
        # and its answer is written that way, bold question then plain reply in one paragraph, and
        # without this the only route to it was a purpose-written script. A line given as a list of
        # {"text": ..., "bold": ...} becomes exactly those runs, in order.
        stuecke = line if isinstance(line, list) else [line]
        for stueck in stuecke:
            run_element = copy.deepcopy(keep)
            end_para_rpr = paragraph._p.find(
                "{http://schemas.openxmlformats.org/drawingml/2006/main}endParaRPr")
            if end_para_rpr is not None:
                # A run appended after endParaRPr breaks the schema's element order; PowerPoint
                # then drops the run silently, no error, while a more tolerant renderer still
                # shows it. TextFrame.clear() keeps endParaRPr on the first paragraph, so this
                # is the common case, not the rare one.
                end_para_rpr.addprevious(run_element)
            else:
                paragraph._p.append(run_element)
            neuer = paragraph.runs[-1]
            roh = str(stueck.get("text", "")) if isinstance(stueck, dict) else str(stueck)
            if isinstance(stueck, dict):
                if "bold" in stueck:
                    neuer.font.bold = bool(stueck["bold"])
                if "italic" in stueck:
                    neuer.font.italic = bool(stueck["italic"])
            _set_run_text(paragraph, neuer, roh, keep)
    return None


_SOFT_BREAK = "\u000b"


def _set_run_text(paragraph, run, text: str, template_r) -> None:
    """Put text into a run, turning a soft line break into the element PowerPoint reads as one.

    A soft break inside a paragraph is U+000B in the text, and `<a:br/>` in the file. Assigning the
    character straight to `run.text` writes the literal `_x000B_` into the XML, and the render then
    shows those seven characters in the middle of the sentence. Found in practice: a
    hand-corrected reference slide carried such a break, and writing its text back put the literal
    on the slide. A newline in the same position means the same thing and is treated the same way.
    """
    teile = str(text).replace("\r\n", "\n").replace("\r", "\n").replace("\n", _SOFT_BREAK)
    stuecke = teile.split(_SOFT_BREAK)
    run.text = stuecke[0]
    letzter = run._r
    for weiterer in stuecke[1:]:
        br = letzter.makeelement(
            "{http://schemas.openxmlformats.org/drawingml/2006/main}br", {})
        letzter.addnext(br)
        neuer_r = copy.deepcopy(template_r)
        br.addnext(neuer_r)
        letzter = neuer_r
        for lauf in paragraph.runs:
            if lauf._r is neuer_r:
                lauf.text = weiterer
                break


def fill_slots(slide, texts: dict, slots: dict, strict: bool) -> list[str]:
    """Put the new text in, addressed by slot name, and report what does not fit.
    Overflow is a refusal rather than a smaller type size: shrinking to fit is how
    a deck stops looking like the slide it was copied from."""
    problems = []
    live = slot_names(slide)
    vergeben: set[int] = set()
    for name, value in texts.items():
        if name not in live:
            problems.append(f"no slot '{name}' on this slide (has: {', '.join(sorted(live))})")
            continue
        lines = value if isinstance(value, list) else [value]
        room = slots.get(name, live[name])["capacity"]

        def _laenge(teil) -> int:
            if isinstance(teil, dict):
                return len(str(teil.get("text", "")))
            if isinstance(teil, list):
                return sum(_laenge(t) for t in teil)
            return len(str(teil))

        length = sum(_laenge(line) for line in lines)
        if length > room:
            problems.append(f"slot '{name}' holds {room} characters, the new text has {length}")
            if strict:
                continue
        frame = _locate_slot_frame(slide, name, live[name]["text"], vergeben)
        if frame is None:
            problems.append(f"slot '{name}' could not be located in the copy")
            continue
        error = _fill_text_frame(frame, lines)
        if error:
            problems.append(f"slot '{name}' {error}")
    return problems


def _resolve_source(source: Path, library: Path) -> Path | None:
    """The harvested deck, found from wherever this is called.

    The index stores the path as it was given at harvest time, which is relative to the directory
    the harvest ran in. Called from anywhere else, that path points at nothing, and the message said
    the deck was "gone", sending someone to look for a file that was never deleted. Found 2026-08-26
    on a real run from the space root.
    """
    for kandidat in (source, library / source, library.parent / source,
                     library.resolve().parent.parent / source):
        if kandidat.is_file():
            return kandidat
    return None


def build(plan_path: Path, out: Path, library: Path, strict: bool) -> int:
    index = json.loads((library / "index.json").read_text(encoding="utf-8"))
    source_deck = _resolve_source(Path(index["source"]), library)
    if source_deck is None:
        roh = Path(index["source"])
        print(f"slide-library: cannot find the harvested deck. The index names {roh}, which is "
              f"relative, and it was tried from here ({Path.cwd()}) and from the library "
              f"({library.resolve()}). Run this from the bundle the harvest ran in, or re-harvest.",
              file=sys.stderr)
        return 2
    plan = json.loads(plan_path.read_text(encoding="utf-8"))
    slides = plan.get("slides") or []
    if not slides:
        print("slide-library: the plan holds no slide", file=sys.stderr)
        return 2

    deck = Presentation(str(source_deck))
    sources = list(deck.slides)
    keep = set(plan.get("keep") or [])
    problems: list[str] = []

    for step, wanted in enumerate(slides, start=1):
        source_id = wanted.get("from")
        entry = next((e for e in index["slides"] if e["id"] == source_id), None)
        if entry is None:
            print(f"slide-library: plan step {step} names slide {source_id}, which is not in the library",
                  file=sys.stderr)
            return 2
        new = clone(deck, sources[source_id - 1])
        found = fill_slots(new, wanted.get("texts") or {}, entry["slots"], strict)
        problems += [f"step {step} ({wanted.get('name', 'unnamed')}): {p}" for p in found]

    # the harvested deck's own slides go, unless the plan keeps some
    xml_slides = deck.slides._sldIdLst
    for position, slide_id in enumerate(list(xml_slides)[:len(sources)], start=1):
        if position in keep:
            continue
        rId = slide_id.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
        deck.part.drop_rel(rId)
        xml_slides.remove(slide_id)

    if not guarded_save(deck, out):
        return 1
    print(f"built {len(slides)} slides into {out}")
    # Read back what was actually written. A reference into nothing makes PowerPoint call the file
    # damaged and strip the shape, while LibreOffice renders it without a word, so neither the render
    # nor any other check here would have caught it.
    lose = dangling_refs(out)
    for eintrag in lose:
        print(f"FAIL: {eintrag}", file=sys.stderr)
    if lose:
        print(f"{len(lose)} reference(s) into nothing. The file will open as damaged in PowerPoint.",
              file=sys.stderr)
        return 1
    if problems:
        print()
        for problem in problems:
            print(f"FAIL: {problem}")
        print(f"\n{len(problems)} open. Shorten the text or pick a slide that carries it.")
        return 1
    return 0


def show(library: Path, slide_id: int | None) -> int:
    index = json.loads((library / "index.json").read_text(encoding="utf-8"))
    for entry in index["slides"]:
        if slide_id and entry["id"] != slide_id:
            continue
        groups = ", ".join(f"{count}x {kind}" for kind, count in entry["groups"].items()) or "-"
        print(f"\nSlide {entry['id']}: {entry['title'] or '(no title)'}  [{entry['layout']}, {groups}]")
        for name, slot in entry["slots"].items():
            print(f"   {name:<16} holds ~{slot['capacity']:>4} chars   now {len(slot['text']):>4}   "
                  f"{slot['size_pt'] or '-'} pt")
    return 0


def check(library: Path, task: str, space_root: Path,
          shape: str | None = None, why: str | None = None) -> int:
    """Print the library, same as `show`, and record that it was looked at for `task`
    (a `workbench/<slug>/` bundle). `library-check-guard` reads this record back, so what it
    proves is that a Compose build happened after a look at the library, not instead of
    one, and it does not judge whether Compose was the right call, only that the cheap
    tiers were on the table when it was made."""
    result = show(library, None)
    # The brand's own approved slides belong in this look, and for a while they were not in it:
    # the guard proved the library had been seen while the cheapest route of all, a slide the user
    # already approved, was not on the list being shown. Measured: a run did the check
    # and then wrote its own build script, with two matching slides sitting unlisted.
    marken = sorted((space_root / "trusted" / "brands").glob("*/slides/INDEX.md"))
    if marken:
        print()
        for index_datei in marken:
            marke = index_datei.parent.parent.name
            eintraege = sorted(f.stem for f in index_datei.parent.glob("*.pptx"))
            print(f"Approved slides of brand {marke!r} ({len(eintraege)}), the cheapest route there "
                  f"is: copy one and fill it.")
            for eintrag in eintraege:
                notiz = index_datei.parent / f"{eintrag}.md"
                zeile = ""
                if notiz.is_file():
                    for z in notiz.read_text(encoding="utf-8").splitlines():
                        if z.startswith("Came from:"):
                            zeile = z
                            break
                print(f"  {eintrag:<28} {zeile}")
            print(f"  -> slide-library.py extract {index_datei.parent}/<name>.pptx --slides 1 "
                  f"--out <new.pptx>, then fill")
    marker_dir = space_root / "zanmai" / "temp" / task
    marker_dir.mkdir(parents=True, exist_ok=True)
    index = json.loads((library / "index.json").read_text(encoding="utf-8"))
    marker_datei = marker_dir / "library-checked.json"
    frueher = {}
    if marker_datei.is_file():
        try:
            frueher = json.loads(marker_datei.read_text(encoding="utf-8"))
        except ValueError:
            frueher = {}
    # Which shape was chosen, and why. Kept per bundle, because the failure this catches is a
    # bundle-level one: the pattern used for the last piece gets used again for a piece whose
    # content has a different shape. Seen on the second piece of a bundle, with the
    # skill's "decide the shape first" already in place: a sentence in a file costs nothing to
    # skip, a line that has to be written and read back does not.
    gewaehlt = list(frueher.get("shapes_chosen") or [])
    if shape:
        schon = [e for e in gewaehlt if e.get("shape") == shape]
        if schon:
            # Said, not refused. Refusing turned "do not always reach for the same one" into "never
            # use the same one twice", and a run then picked a different pattern for every piece
            # whether or not it carried the content. A product family looks alike; that is not a
            # fault. What this is for is that the second use is a decision rather than a reflex.
            print(f"\n{shape!r} is already used in this bundle, {len(schon)}x, most recently for "
                  f"{schon[-1].get('for', '?')}: {schon[-1].get('why', 'no reason recorded')}\n"
                  f"Fine where this content has the same shape. Worth a second look where it does "
                  f"not, because reaching for the last one used is the cheap move, not the right one.")
        gewaehlt.append({"shape": shape, "why": why or "", "for": task,
                         "at": datetime.now(timezone.utc).isoformat(timespec="seconds")})
    marker = {
        "checked_at": datetime.now(timezone.utc).isoformat(timespec="seconds"),
        "library": str(library),
        "slides_seen": len(index["slides"]),
        "shapes_chosen": gewaehlt,
    }
    marker_datei.write_text(json.dumps(marker, indent=2, ensure_ascii=False) + "\n", encoding="utf-8")
    print(f"\nrecorded: {task} looked at {marker['slides_seen']} slide(s) in {library}")
    if gewaehlt:
        print("shapes chosen in this bundle so far:")
        for eintrag in gewaehlt:
            print(f"  {eintrag['shape']:<24} {eintrag['why'][:70]}")
    return result


def _shape_paths(shapes, prefix="", parents=()):
    """(path, shape, parents) for every shape in the tree, depth-first. `path` is the
    slash-joined chain of names down to this shape, which is what `--shape` accepts when
    a bare name repeats across sibling groups (a template's repeated cards do, verified on
    a real deck where every card's text box was named "Rectangle 63"). `parents` is the
    chain of enclosing group shapes, outermost first, which is what a caller needs to
    convert a slide-space distance into this shape's own coordinate units."""
    for sh in shapes:
        path = f"{prefix}/{sh.name}"
        yield path, sh, parents
        if sh.shape_type == 6:
            yield from _shape_paths(sh.shapes, path, parents + (sh,))


def _find_shape(deck, wanted: str, slide_no: int | None):
    """Every (slide_index, path, shape, parents) matching `wanted`, by full path first,
    then by bare name. Ambiguous or missing is the caller's problem to report, not a guess
    made here: picking the wrong one of several identically named cards is worse than an
    error that names every candidate."""
    slides = [(slide_no, deck.slides[slide_no - 1])] if slide_no else list(enumerate(deck.slides, start=1))
    hits = []
    for index, slide in slides:
        candidates = list(_shape_paths(slide.shapes))
        exact = [(index, p, sh, parents) for p, sh, parents in candidates if p in (wanted, "/" + wanted)]
        hits += exact or [(index, p, sh, parents) for p, sh, parents in candidates if sh.name == wanted]
    return hits


def _group_scale(group) -> tuple[float, float]:
    """(sx, sy): how one unit in this group's own child space maps to one unit in the
    space its own xfrm is expressed in (its parent group's child space, or slide space for
    a top-level group)."""
    xfrm = group._element.grpSpPr.xfrm
    ext, chExt = xfrm.find(qn("a:ext")), xfrm.find(qn("a:chExt"))
    cx, cy = int(chExt.get("cx")), int(chExt.get("cy"))
    return (int(ext.get("cx")) / cx if cx else 1.0,
            int(ext.get("cy")) / cy if cy else 1.0)


def _regroup(group) -> None:
    """After a descendant moved, keep this group's frame exactly containing its direct
    children: recompute chOff/chExt as their union in this group's own child space, and
    grow off/ext by the same amount so the scale to its own parent stays fixed. Without
    this a child that moved outside the old frame is silently clipped or rescaled by
    PowerPoint, which trusts chOff/chExt rather than re-measuring its children."""
    kids = [k for k in group.shapes if k.left is not None]
    if not kids:
        return
    new_x = min(k.left for k in kids)
    new_y = min(k.top for k in kids)
    new_w = max(k.left + k.width for k in kids) - new_x
    new_h = max(k.top + k.height for k in kids) - new_y

    xfrm = group._element.grpSpPr.xfrm
    off, ext = xfrm.find(qn("a:off")), xfrm.find(qn("a:ext"))
    chOff, chExt = xfrm.find(qn("a:chOff")), xfrm.find(qn("a:chExt"))
    sx, sy = _group_scale(group)

    off.set("x", str(int(off.get("x")) + round((new_x - int(chOff.get("x"))) * sx)))
    off.set("y", str(int(off.get("y")) + round((new_y - int(chOff.get("y"))) * sy)))
    ext.set("cx", str(round(new_w * sx)))
    ext.set("cy", str(round(new_h * sy)))
    chOff.set("x", str(new_x))
    chOff.set("y", str(new_y))
    chExt.set("cx", str(new_w))
    chExt.set("cy", str(new_h))


def swap_image(deck_path: Path, wanted: str, quelle: str, slide_no: int | None, out: Path) -> int:
    """Replace one picture with another, whole, and keep the target's position.

    Swapping only the image source is the obvious move and it is wrong. Every picture carries its
    own `a:srcRect`, a crop cut for exactly its own artwork: a set of icons can share one viewBox
    and still hold different amounts of whitespace inside it, and the crop is what evens that out.
    Leave the old crop on the new artwork and the icon slides out of its field and gets cut off --
    visible in a render, invisible to every check here, because the shape's box never changed.
    Seen in practice, after the naive version had been built first and the render
    disproved it.

    So the whole `p:pic` comes across, crop included, and only the placement is set afterwards:
    height and left edge from the shape being replaced, width from the source's own aspect ratio.
    `--from` is `<deck>:<slide>:<shape>`, and where that deck is this deck, the source slide has to
    still be in the file: once it is dropped, its image parts are not reachable any more.
    """
    teile = quelle.split(":")
    if len(teile) != 3:
        print(f"slide-library: --from takes <deck>:<slide>:<shape>, got {quelle!r}", file=sys.stderr)
        return 2
    quell_deck_pfad, quell_folie, quell_name = Path(teile[0]), teile[1], teile[2]
    if not quell_deck_pfad.is_file():
        print(f"slide-library: no deck at {quell_deck_pfad}", file=sys.stderr)
        return 2
    try:
        quell_folie_nr = int(quell_folie)
    except ValueError:
        print(f"slide-library: --from needs a slide number, got {quell_folie!r}", file=sys.stderr)
        return 2

    deck = Presentation(str(deck_path))
    ziel_treffer = _find_shape(deck, wanted, slide_no)
    if len(ziel_treffer) != 1:
        if not ziel_treffer:
            print(f"slide-library: no shape '{wanted}' in {deck_path.name}", file=sys.stderr)
        else:
            print(f"slide-library: '{wanted}' is not unique, qualify with its full path:", file=sys.stderr)
            for index, path, _, _ in ziel_treffer:
                print(f"  slide {index}: {path}", file=sys.stderr)
        return 2
    ziel_index, _pfad, ziel_shape, _eltern = ziel_treffer[0]

    quell_deck = deck if quell_deck_pfad.resolve() == deck_path.resolve() else Presentation(str(quell_deck_pfad))
    quell_treffer = [h for h in _find_shape(quell_deck, quell_name, quell_folie_nr)]
    if len(quell_treffer) != 1:
        print(f"slide-library: '{quell_name}' on slide {quell_folie_nr} of {quell_deck_pfad.name} "
              f"matched {len(quell_treffer)} shape(s)", file=sys.stderr)
        return 2
    _qi, _qp, quell_shape, _qe = quell_treffer[0]
    if quell_shape._element.tag != qn("p:pic"):
        print(f"slide-library: '{quell_name}' is not a picture", file=sys.stderr)
        return 2

    quell_folie_obj = quell_deck.slides[quell_folie_nr - 1]
    ziel_folie_obj = deck.slides[ziel_index - 1]

    kopie = copy.deepcopy(quell_shape._element)
    _carry_relationships(ziel_folie_obj, quell_folie_obj, kopie)

    # Placement: the target's height and left edge, the source's own aspect ratio for the width.
    hoehe = ziel_shape.height
    q_breite, q_hoehe = quell_shape.width, quell_shape.height
    breite = int(q_breite * (hoehe / q_hoehe)) if q_hoehe else ziel_shape.width
    ziel_shape._element.addnext(kopie)
    eltern = ziel_shape._element.getparent()
    eltern.remove(ziel_shape._element)
    xfrm = kopie.find(qn("p:spPr")).find(qn("a:xfrm"))
    off, ext = xfrm.find(qn("a:off")), xfrm.find(qn("a:ext"))
    off.set("x", str(ziel_shape.left))
    off.set("y", str(ziel_shape.top))
    ext.set("cx", str(breite))
    ext.set("cy", str(hoehe))

    if not guarded_save(deck, out):
        return 1
    crop = kopie.find(qn("p:blipFill")) is not None and \
        kopie.find(qn("p:blipFill")).find(qn("a:srcRect")) is not None
    print(f"swapped '{wanted}' on slide {ziel_index} for '{quell_name}' -> {out}")
    print(f"  the whole picture came across{', crop included' if crop else ''}; placement is the "
          f"old shape's height and left edge, width from the source's aspect ratio")
    lose = dangling_refs(out)
    for eintrag in lose:
        print(f"FAIL: {eintrag}", file=sys.stderr)
    return 1 if lose else 0


def nudge(deck_path: Path, wanted: str, dx_in: float, dy_in: float, slide_no: int | None, out: Path) -> int:
    deck = Presentation(str(deck_path))
    hits = _find_shape(deck, wanted, slide_no)
    if not hits:
        print(f"slide-library: no shape '{wanted}' on {'slide ' + str(slide_no) if slide_no else 'any slide'}",
              file=sys.stderr)
        return 2
    if len(hits) > 1:
        print(f"slide-library: '{wanted}' is not unique, qualify with its full path:", file=sys.stderr)
        for index, path, _, _ in hits:
            print(f"  slide {index}: {path}", file=sys.stderr)
        return 2

    index, path, shape, parents = hits[0]
    sx, sy = 1.0, 1.0
    for group in parents:
        gsx, gsy = _group_scale(group)
        sx *= gsx
        sy *= gsy
    dx_emu = round(dx_in * EMU_PER_INCH / sx)
    dy_emu = round(dy_in * EMU_PER_INCH / sy)

    left, top, width, height = shape.left, shape.top, shape.width, shape.height
    shape.left, shape.top, shape.width, shape.height = left + dx_emu, top + dy_emu, width, height
    for group in reversed(parents):
        _regroup(group)

    if not guarded_save(deck, out):
        return 1
    print(f"slide {index}: moved {path} by {dx_in:+.3f}in / {dy_in:+.3f}in -> {out}")
    return 0


CHAR_WIDTH_EM = 0.52
# Average glyph advance as a fraction of point size, calibrated against the reference
# build's bold numerals, the case that actually painted past its own box. Coarse like
# CHAR_DENSITY above, and for the same reason: it only has to tell "still fits" from
# "paints past its box", not predict a layout engine exactly. Used only where the real
# font file below cannot be found -- a fallback, not the primary measurement.


def _painted_width(char_count: int, size_pt: float, bold: bool) -> int:
    per_char = (size_pt / 72.0) * EMU_PER_INCH * (CHAR_WIDTH_EM + (0.05 if bold else 0.0))
    return round(char_count * per_char)


def _font_file(family: str, bold: bool) -> Path | None:
    """A real font file for `family`, matched by filename in the standard font folders
    (same list `installed_font_families()` in design-check.py scans, one level further:
    a file, not just a name), closest weight to `bold` preferred. None means no match on
    this machine; callers fall back to the calibrated estimate rather than fail."""
    roots = [Path("/System/Library/Fonts"), Path("/Library/Fonts"), Path.home() / "Library/Fonts",
             Path("/usr/share/fonts"), Path("/usr/local/share/fonts"), Path.home() / ".fonts"]
    fam_key = re.sub(r"[^a-z0-9]", "", family.lower()) if family else ""
    if not fam_key:
        return None
    candidates = []
    for root in roots:
        if not root.is_dir():
            continue
        for file in root.rglob("*"):
            if file.suffix.lower() not in (".ttf", ".otf", ".ttc"):
                continue
            if re.sub(r"[^a-z0-9]", "", file.stem.lower()).startswith(fam_key):
                candidates.append(file)
    if not candidates:
        return None
    variant_markers = ("italic", "narrow", "condensed", "rounded", "unicode", "oblique")
    plain = [f for f in candidates if not any(m in f.stem.lower() for m in variant_markers)] or candidates
    if bold:
        # A literal "bold" not preceded by "extra"/"semi" is the correct match for a plain bold
        # request. "Black" is not: on this machine "Arial Black" is a distinct shipped family,
        # not a weight of "Arial", so treating a name match on "black" as "Arial, but bold" picked
        # the wrong file entirely -- found while building this function, not assumed away.
        exact = [f for f in plain if re.search(r"(?<!extra)(?<!semi)bold", f.stem.lower())]
        if exact:
            return min(exact, key=lambda f: len(f.stem))
        for keyword in ("black", "heavy", "extrabold", "semibold"):
            matches = [f for f in plain if keyword in f.stem.lower()]
            if matches:
                return min(matches, key=lambda f: len(f.stem))
    weight_markers = ("bold", "black", "light", "thin", "medium", "semibold", "extrabold", "heavy")
    unweighted = [f for f in plain if not any(m in f.stem.lower() for m in weight_markers)]
    return min(unweighted or plain, key=lambda f: len(f.stem))


def _real_ink_bbox(text: str, family: str | None, size_pt: float, bold: bool):
    """(left_bearing_in, width_in) of `text` set in this exact font file at this exact
    size, not an average -- the same "measured, not eyeballed" standard `capacity()`
    already applies to a slot's character count, extended to where the ink itself sits.
    None if Pillow is missing, the family cannot be resolved to a file, or the file
    fails to parse: an unmeasured run is left to the calibrated fallback, never guessed
    into a false positive or a false clean."""
    if not family or not text:
        return None
    try:
        from PIL import ImageFont
    except ImportError:
        return None
    path = _font_file(family, bold)
    if path is None:
        return None
    # Rendered rather than asked for, and rendered large. `getbbox` reports the ink box measured
    # from the pen position, so its left value is zero for every string in every face: the side
    # bearing it is supposed to carry has already been folded into the layout. Measuring the drawn
    # mask is what actually finds the bearing, and drawing at eight times the size is what stops
    # the answer from snapping to whole pixels, which at 11 pt is a third of a millimetre of
    # quantisation on a quantity of the same order.
    #
    # This is why a title and its body could sit visibly apart while the check called them level:
    # it was comparing two zeroes. On a page with six sizes down one edge that is the difference
    # anybody sees at a glance.
    ueber = 8
    try:
        font = ImageFont.truetype(str(path), size=max(1, round(size_pt * ueber)))
        layout_left, _t, layout_right, _b = font.getbbox(text)
        maske = font.getmask(text).getbbox()
    except Exception:  # noqa: BLE001 -- a font file that fails to parse falls back, it does not crash the check
        return None
    # Two measurements, because each one is blind to the case the other catches. `getbbox` gives
    # the box the layout occupies, widened where ink hangs past the pen on the left, so it reports
    # a negative for a letter that overhangs and a flat zero for one that simply starts inset.
    # The rendered mask gives that inset, but it is clipped at zero and cannot see an overhang.
    # A face is one or the other per string, never both, so whichever is non-zero is the answer.
    if layout_left < 0:
        left, right = layout_left, layout_right
    elif maske is not None:
        left, right = maske[0], maske[2]
    else:                   # a string that paints nothing, a space for instance
        return None
    return (left / ueber) / 72.0, ((right - left) / ueber) / 72.0


def _leaf_rects(shapes, tf=(1.0, 0.0, 1.0, 0.0), mit_text=True):
    """(shape, (left, top, width, height) in slide EMU) for every leaf, groups resolved through
    their own scale and offset so a card's children compare against a sibling card's in one shared
    coordinate space, not each in its own group's raw units.

    `mit_text=True` keeps only leaves that carry text, which is what every check that measures
    wording wants. Pass False where the question is about the surface rather than the words: a card
    is often a filled shape with no text at all, and the words sit in separate boxes over it, so a
    walk that skips empty shapes cannot see the card.
    """
    ax, bx, ay, by = tf
    out = []
    for sh in shapes:
        if sh.shape_type == 6:
            sx, sy = _group_scale(sh)
            xfrm = sh._element.grpSpPr.xfrm
            off, chOff = xfrm.find(qn("a:off")), xfrm.find(qn("a:chOff"))
            nested = (ax * sx, ax * (int(off.get("x")) - int(chOff.get("x")) * sx) + bx,
                      ay * sy, ay * (int(off.get("y")) - int(chOff.get("y")) * sy) + by)
            out += _leaf_rects(sh.shapes, nested, mit_text)
            continue
        if sh.left is None or sh.width is None or sh.height is None:
            continue
        if mit_text and (not sh.has_text_frame or not text_of(sh)):
            continue
        out.append((sh, (ax * sh.left + bx, ay * sh.top + by, ax * sh.width, ay * sh.height)))
    return out


def _painted_rect(shape, rect):
    """Widen `rect` to the glyph's real ink where `wrap="none"` lets the run paint past its
    own box. Seen in practice: an 80pt digit box sized for two characters still
    painted a three-character run, and the saved bounding box said nothing overlapped."""
    body_pr = shape.text_frame._txBody.find(qn("a:bodyPr"))
    if body_pr is None or body_pr.get("wrap") != "none":
        return rect
    left, top, width, height = rect
    widest = width
    for paragraph in shape.text_frame.paragraphs:
        length = sum(len(r.text) for r in paragraph.runs)
        if not length:
            continue
        size = max((r.font.size.pt for r in paragraph.runs if r.font.size), default=12.0)
        bold = any(r.font.bold for r in paragraph.runs)
        family = next((r.font.name for r in paragraph.runs if r.font.name), None)
        run_text = "".join(r.text for r in paragraph.runs)
        real = _real_ink_bbox(run_text, family, size, bold)
        if real is not None:
            widest = max(widest, round(real[1] * EMU_PER_INCH))
        else:
            widest = max(widest, _painted_width(length, size, bold))
    return (left, top, widest, height)


def _ink_left_offset(shape) -> int | None:
    """How far this shape's first line actually starts painting from its own box's
    left edge, in EMU. None where that cannot be measured (no run, no named font, no
    matching file on this machine) -- such a shape is left out of `align_check` rather
    than assumed aligned or assumed not.

    Three things sit between the box edge and the first glyph, and for a long time only the third
    was counted: the frame's own left inset, the paragraph's indent, and the glyph's side bearing.
    The inset is the one that bites, because it is usually inherited rather than written: a title
    taking the layout's default 0.1 inch over a lead paragraph set to zero reads as misaligned on
    every slide, while both boxes sit on exactly the same edge. Seen twice, the
    second time on 13 of 16 slides, with the check reporting 47 groups and no fault.
    """
    rahmen = shape.text_frame
    inset = rahmen.margin_left if rahmen.margin_left is not None else 91440
    for paragraph in rahmen.paragraphs:
        if not paragraph.runs:
            continue
        run = paragraph.runs[0]
        if not run.text.strip() or not run.font.size:
            return None
        einzug = 0
        pPr = paragraph._pPr
        if pPr is not None:
            for name in ("marL", "indent"):
                wert = pPr.get(name)
                if wert:
                    einzug += int(wert)
        real = _real_ink_bbox(run.text, run.font.name, run.font.size.pt, bool(run.font.bold))
        if real is None:
            return None
        return round(real[0] * EMU_PER_INCH) + inset + einzug
    return None


def align_check(deck_path: Path, slide_no: int | None) -> int:
    """Two text frames whose boxes sit on the same left edge can still read as
    misaligned: a large bold face and a small regular face rarely share the same left
    side bearing, so their ink starts at different points even though the boxes agree
    exactly (found on a real deck, a title and its body text). Boxes within 0.02in are
    treated as "meant to align"; among those, ink-left positions more than 0.02in apart
    are reported. A shape whose ink cannot be measured is left out of its group, never
    silently counted as aligned."""
    deck = Presentation(str(deck_path))
    slides = [(slide_no, deck.slides[slide_no - 1])] if slide_no else list(enumerate(deck.slides, start=1))
    # Half a millimetre is invisible on a wall and plain on a sheet held at arm's length, so the
    # tolerance follows the page like the layout floors do. Grouping stays at the coarse value:
    # that answers "were these meant to line up", which is a question about intent, while the
    # comparison answers "do they", which is a question about the eye.
    druck = _ist_druckstueck(inches(deck.slide_width), inches(deck.slide_height))
    gruppen_tol = round(0.02 * EMU_PER_INCH)
    tolerance = round((ALIGN_TOL_PRINT if druck else ALIGN_TOL_SCREEN) * EMU_PER_INCH)
    problems = []
    groups_checked = 0
    for index, slide in slides:
        leaves = _leaf_rects(slide.shapes)
        leaves = sorted(leaves, key=lambda pair: pair[1][0])
        groups: list[list] = []
        for sh, rect in leaves:
            for group in groups:
                if abs(group[0][1][0] - rect[0]) <= gruppen_tol:
                    group.append((sh, rect))
                    break
            else:
                groups.append([(sh, rect)])
        for group in groups:
            if len(group) < 2:
                continue
            measured = [(sh, rect[0] + offset) for sh, rect in group
                        if (offset := _ink_left_offset(sh)) is not None]
            if len(measured) < 2:
                continue
            groups_checked += 1
            for i in range(len(measured)):
                sh_a, left_a = measured[i]
                for j in range(i + 1, len(measured)):
                    sh_b, left_b = measured[j]
                    if abs(left_a - left_b) > tolerance:
                        # The number to move by, not just the number that is wrong. Whoever fixes
                        # this has to shift one box left by exactly this much, so saying how far
                        # apart they are and leaving the direction to be worked out again is half
                        # an answer.
                        weiter, zurueck = ((sh_a, sh_b) if left_a > left_b else (sh_b, sh_a))
                        versatz = Emu(abs(left_a - left_b)).inches
                        problems.append(
                            "slide %d: '%s' and '%s' share a box-left edge but their ink starts "
                            "%.3f inch apart: move '%s' left by %.3f inch (%.2f mm)"
                            % (index, text_of(sh_a)[:20], text_of(sh_b)[:20], versatz,
                               text_of(weiter)[:20], versatz, versatz * 25.4))
    for problem in problems:
        print(f"FAIL: {problem}")
    art = "print" if druck else "screen"
    print(f"{groups_checked} shared-edge group(s) checked, {len(problems)} misalignment(s) "
          f"[{art}, ink apart over {Emu(tolerance).inches:.3f}in]"
          + _nichts_geprueft(groups_checked))
    return 1 if problems else 0


# What page furniture may drift between pages. Nothing, in effect: a footer, a logo, a page number
# or a contact strip that moves a quarter of a millimetre between two pages was never meant to, and
# where it genuinely was, the piece says so by not repeating the element. Every other check in this
# file reads one page at a time, so this whole class of fault went past all of them green.
FURNITURE_TOL = 0.01         # inch; a hard floor, because drift here is never intentional
_GENERIC_NAME = re.compile(r"^(TextBox|Rectangle|Oval|Picture|Freeform|Group|Content Placeholder|"
                           r"Title|Subtitle|Text Placeholder)\s*\d*$", re.IGNORECASE)


# How much of the page counts as its edge. Page furniture lives there and content does not: a
# kicker and a logo at the top, a footer, a page number and a contact strip at the bottom. Text
# repeated in the middle of a page is a sentence the piece makes twice, which is a normal thing for
# a piece to do and none of this check's business. Without this, a flyer that carries its core
# sentence on both pages had it reported three times as furniture that had moved.
FURNITURE_ZONE = 0.125       # top and bottom eighth of the page
# How far apart two things may sit and still be the same slot on two pages. A quarter inch: wide
# enough that a slot which slipped is recognised as that slot and then reported, narrow enough that
# the next piece of furniture down is never taken for it. Paired loosely, compared tightly, and the
# gap between two slots is what sets the ceiling: a kicker and the title under it sit half an inch
# apart, so half an inch as the window made every such pair ambiguous and stopped the pairing.
_ZONEN_FENSTER = round(0.25 * EMU_PER_INCH)


def _in_furniture_zone(rect, seitenhoehe: int, zone: float) -> bool:
    """Whether this shape sits wholly in the top or bottom band of the page."""
    oben, hoehe = rect[1], rect[3]
    band = seitenhoehe * zone
    return (oben + hoehe) <= band or oben >= (seitenhoehe - band)


def _furniture_key(shape) -> str | None:
    """What makes this shape the same element as one on another page, or None.

    The text first, because a footer carries the same words on every page and that is the thing a
    reader recognises it by. The shape name second, and only where the builder gave it one: the
    default names carry a counter that restarts per page, so `TextBox 3` on page one and
    `TextBox 3` on page two are not the same element and matching on them would invent faults.
    """
    text = text_of(shape).strip() if shape.has_text_frame else ""
    if text:
        return "text:" + re.sub(r"\s+", " ", text.lower())[:80]
    name = (shape.name or "").strip()
    if name and not _GENERIC_NAME.match(name):
        return "name:" + name.lower()
    return None


def furniture_check(deck_path: Path, tol_in: float = None, zone: float = None) -> int:
    """An element that repeats across pages, sitting somewhere else on one of them.

    A footer at the same distance from the bottom on page one and two millimetres higher on page
    two passes `align-check`, `overlap-check`, `overflow-check` and every render, because each of
    those reads a single page and on each single page nothing is wrong. What is wrong is only
    visible with both pages side by side, which until now was a person's job on every piece.

    An element appearing twice on the same page is skipped rather than guessed at: two cards
    carrying the same word are not one element in two places.
    """
    tol = FURNITURE_TOL if tol_in is None else tol_in
    band = FURNITURE_ZONE if zone is None else zone
    deck = Presentation(str(deck_path))
    seiten = list(enumerate(deck.slides, start=1))
    if len(seiten) < 2:
        print("1 page, nothing to compare across pages")
        return 0

    vorkommen: dict[str, dict[int, tuple]] = {}
    mehrfach: dict[str, set[int]] = {}
    nach_zone: dict[str, dict[int, list]] = {}
    for index, slide in seiten:
        # A logo is page furniture and carries no text, so a walk over text-bearing shapes
        # could never pair one: the element most likely to wander between pages was the
        # one this check could not see.
        for shape, rect in _leaf_rects(slide.shapes, mit_text=False):
            if shape.left is None or shape.top is None:
                continue
            if not _in_furniture_zone(rect, deck.slide_height, band):
                continue
            key = _furniture_key(shape)
            # The label is cut to 24 characters for reading; the full key travels with it. Deciding
            # anything on the cut version compares two prefixes and calls them identical: two
            # kickers reading "NIS-2 UND CYBER-COMPLIANCE" and "NIS-2 UND CYBER-COMPLIANCE ·
            # MANAGED FIREWALL" are the same string for the first 24 characters, so the height pass
            # concluded the wording had not changed and left the pair to a content pass that could
            # not take it either. The same file with a shorter kicker paired correctly, which is
            # what made it look like a geometry problem.
            eintrag = (rect[0], rect[1], rect[2], rect[3],
                       text_of(shape)[:24] if shape.has_text_frame else shape.name,
                       key or f"shape:{shape.name}")
            # A second key, for the same slot carrying different words: a kicker reading
            # "Campaign" on page one and "Campaign, Product" on page two is one piece of
            # furniture, and matching on the words cannot see that. Collected by zone here,
            # paired by height below.
            zone_name = "top" if (rect[1] + rect[3]) <= deck.slide_height * band else "bottom"
            nach_zone.setdefault(zone_name, {}).setdefault(index, []).append((rect, eintrag))
            if key is None:
                continue
            if index in vorkommen.get(key, {}):
                mehrfach.setdefault(key, set()).add(index)
                continue
            vorkommen.setdefault(key, {})[index] = eintrag

    problems, geprueft = [], 0
    gepaart: set[str] = set()
    grenze = round(tol * EMU_PER_INCH)
    for key, je_seite in sorted(vorkommen.items()):
        seiten_mit = {s for s in je_seite if s not in mehrfach.get(key, set())}
        if len(seiten_mit) < 2:
            continue
        geprueft += 1
        basis_seite = min(seiten_mit)
        links, oben, breite, hoehe, label, _schluessel = je_seite[basis_seite]
        gepaart.add(f"{label} ({len(seiten_mit)} pages)")
        for seite in sorted(seiten_mit - {basis_seite}):
            l2, o2, b2, h2, _label2, _key2 = je_seite[seite]
            for wert_a, wert_b, was in ((links, l2, "left"), (oben, o2, "top"),
                                        (breite, b2, "width"), (hoehe, h2, "height")):
                if abs(wert_a - wert_b) > grenze:
                    problems.append(
                        "'%s' sits at a different %s on page %d than on page %d: %.3f inch apart "
                        "(move it to %.3f in)"
                        % (label, was, seite, basis_seite, Emu(abs(wert_a - wert_b)).inches,
                           Emu(wert_a).inches))
    # The height pass, for what the content pass could not pair. Nearest neighbour inside the zone
    # rather than a grid: a bucket edge is arbitrary, and two slots either side of one are never
    # compared at all. That is what the column key kept producing, in both directions. The frame's
    # left edge separated slots that belong together, because alignment shifts every frame by its
    # own first glyph's side bearing; measuring the ink edge instead then merged everything,
    # because a piece that aligns its ink properly has all of its left-aligned furniture on exactly
    # one line. A key that gets worse the more carefully a piece is set is the wrong key.
    #
    # Height works because furniture is laid out by height: the kicker at the top, the footer at
    # the bottom, and nothing else at either. Paired loosely and compared tightly, so a slot that
    # slipped a tenth of an inch is still recognised as that slot and then reported for slipping.
    for zone_name, je_seite in sorted(nach_zone.items()):
        if len(je_seite) < 2:
            continue
        basis_seite = min(je_seite)
        for rect, eintrag in je_seite[basis_seite]:
            oben, breite = rect[1], rect[2]
            partner: dict[int, tuple] = {}
            eindeutig = True
            for seite, kandidaten in je_seite.items():
                if seite == basis_seite:
                    continue
                nah = [(r, e) for r, e in kandidaten
                       if abs(r[1] - oben) <= _ZONEN_FENSTER
                       and abs(r[2] - breite) <= _ZONEN_FENSTER]
                if len(nah) > 1:
                    eindeutig = False   # two candidates at this height: not readable, not guessed
                    break
                if nah:
                    partner[seite] = nah[0]
            if not eindeutig or not partner:
                continue
            label, schluessel = eintrag[4], eintrag[5]
            if all(e[5] == schluessel for _r, e in partner.values()):
                continue      # same wording throughout: the content pass already had this one
            geprueft += 1
            gepaart.add(f"{label} and the same slot with other wording ({len(partner) + 1} pages)")
            for seite in sorted(partner):
                r2, e2 = partner[seite]
                if abs(oben - r2[1]) > grenze:
                    problems.append(
                        "'%s' on page %d and '%s' on page %d are the same slot and sit %.3f inch "
                        "apart vertically (move it to %.3f in)"
                        % (e2[4], seite, label, basis_seite, Emu(abs(oben - r2[1])).inches,
                           Emu(oben).inches))
    for problem in problems:
        print(f"FAIL: {problem}")
    # Named, not counted. "1 element checked" reads like full coverage and is not checkable against
    # any expectation: on a real flyer the kicker was silently left out, because it carried the
    # product name on page two and the strings stopped matching. Whoever reads this can only notice
    # that if the list says what was paired.
    for label in sorted(gepaart):
        print(f"  paired across pages: {label}")
    print(f"{geprueft} repeating element(s) checked across {len(seiten)} page(s), "
          f"{len(problems)} drift(s) [top and bottom {band:g} of the page, over {tol:g}in]"
          + _nichts_geprueft(geprueft))
    return 1 if problems else 0


_FOLIENPRAEFIX = re.compile(r"^slide \d+: ")


def _ohne_folie(befund: str) -> str:
    """A finding without its slide number, so the same finding matches across two files.

    A slide lifted out of a source deck keeps its shapes and their text but almost never its
    number, so comparing whole lines would report every inherited finding as new.
    """
    return _FOLIENPRAEFIX.sub("", befund)


def _nur_neue(problems: list[str], baseline: Path | None, pruefer) -> tuple[list[str], int]:
    """The findings that are not already in `baseline`, and how many were dropped.

    Why this exists: on a corporate template `overlap-check` reported fifteen overlaps on a built
    file and the identical fifteen on the untouched original, because that template lays its bands
    over each other on purpose. Without this, the next run repairs fifteen things that were built
    that way, which is the most expensive kind of mistake, correcting something correct.
    """
    if baseline is None:
        return problems, 0
    bekannt = {_ohne_folie(b) for b in pruefer(baseline)}
    neu = [p for p in problems if _ohne_folie(p) not in bekannt]
    return neu, len(problems) - len(neu)


def overlap_check(deck_path: Path, slide_no: int | None, baseline: Path | None = None,
                  als_liste: bool = False):
    deck = Presentation(str(deck_path))
    slides = [(slide_no, deck.slides[slide_no - 1])] if slide_no else list(enumerate(deck.slides, start=1))
    problems = []
    pairs = 0
    for index, slide in slides:
        leaves = [(sh, _painted_rect(sh, rect)) for sh, rect in _leaf_rects(slide.shapes)]
        for i in range(len(leaves)):
            sh_a, (la, ta, wa, ha) = leaves[i]
            for j in range(i + 1, len(leaves)):
                sh_b, (lb, tb, wb, hb) = leaves[j]
                pairs += 1
                ox = min(la + wa, lb + wb) - max(la, lb)
                oy = min(ta + ha, tb + hb) - max(ta, tb)
                if ox > 0 and oy > 0:
                    problems.append(
                        "slide %d: '%s' over '%s', %.3f x %.3f inch"
                        % (index, text_of(sh_a)[:24], text_of(sh_b)[:24], Emu(ox).inches, Emu(oy).inches))
    if als_liste:
        return problems
    problems, geerbt = _nur_neue(
        problems, baseline, lambda p: overlap_check(p, None, als_liste=True))
    for problem in problems:
        print(f"FAIL: {problem}")
    print(f"{pairs} pair(s) checked on {len(slides)} slide(s), {len(problems)} overlap(s)"
          + (f", {geerbt} already in {baseline.name}" if baseline else "")
          + _nichts_geprueft(pairs))
    return 1 if problems else 0


_THEME_PLACEHOLDER = {"+mn-lt": "minor", "+mj-lt": "major", "+mn-ea": "minor", "+mj-ea": "major",
                      "+mn-cs": "minor", "+mj-cs": "major"}


def _cell_family(cell, theme: dict | None) -> str | None:
    """The typeface a table cell really renders in: what its own runs name, resolved through
    the theme placeholders, and the theme's body face where a run names nothing."""
    for paragraph in cell.text_frame.paragraphs:
        for run in paragraph.runs:
            family = _resolve_family(run.font.name, theme)
            if family:
                return family
    return (theme or {}).get("minor")


LAYOUT_MIN_MARGIN = 0.5      # inch from the slide edge; a deck may hold itself to more
LAYOUT_MIN_PT = 12.0         # the lowest size still readable on a shared screen
LAYOUT_MIN_PAD = 0.08        # inch between text and the edge of the filled box it sits in

# The same three floors for a page somebody holds in their hand. A shared screen is read from four
# metres away and a sheet of paper from forty centimetres, so the screen numbers applied to print
# turn a check into something that has to be waved through: 10.5 pt body text and a 9 pt caption
# are ordinary on A4 and would be reported as faults on every single page. A check that gets
# overridden as routine has stopped checking, and the run where it is right goes past unnoticed
# with all the others.
#
# The margin is not simply smaller: an office laser cannot print to the edge and clips roughly a
# third of an inch, so this floor is about the machine rather than about taste.
PRINT_MIN_MARGIN = 0.35      # inch from the page edge, near what a desktop printer refuses to reach
PRINT_MIN_PT = 8.0           # a caption, a footnote, a legal line; under this nobody meant it
PRINT_MIN_PAD = 0.05         # inch between text and the edge of the box it sits in

# How far two ink edges may sit apart before somebody sees it. A slide is read across a room, a
# sheet at arm's length, and between a 24 pt bold and a 10.5 pt regular the side bearings differ by
# about a quarter of a millimetre: invisible on the wall, plain on paper with six sizes down one
# edge, and under the screen tolerance either way, which is why it was never reported.
ALIGN_TOL_SCREEN = 0.02
ALIGN_TOL_PRINT = 0.006

# Paper, in inches, both ways round. A deck is 13.33 x 7.5 or 10 x 7.5 and matches none of these.
_PAPIER = ((8.268, 11.693), (5.827, 8.268), (8.5, 11.0), (11.693, 16.535), (7.165, 10.118))


def _ist_druckstueck(breite: float, hoehe: float) -> bool:
    """Whether these page dimensions are a sheet of paper rather than a slide.

    Read off the size itself, because that is what is true regardless of how the file was made:
    the same PowerPoint writes both, and nothing else in the file says which one was meant.
    """
    for w, h in _PAPIER:
        for a, b in ((w, h), (h, w)):
            if abs(breite - a) < 0.12 and abs(hoehe - b) < 0.12:
                return True
    return False


def layout_check(deck_path: Path, slide_no: int | None = None, min_margin: float = None,
                 min_pt: float = None, min_pad: float = None) -> int:
    """Three faults that every other check here is blind to, because they are about where a
    shape sits rather than what is in it.

    `overflow-check` reads the box; this reads where the box is. A card that runs past the
    slide edge, a caption at 8 pt, a heading whose text touches the edge of its own fill:
    all three paint, none of them wrap, and none of them is an overlap. Found 2026-08-26
    while building the wireframe library, where the first run reported 440 of them across
    21 slides and every single one had passed the existing checks.

    The floors are deliberately low. This says "nobody meant that", not "this is the house
    style": a deck that holds itself to wider margins passes its own numbers in.

    **Which floors, though, follows from the page.** A4 or A5 gets the print set, anything else the
    screen set, because a sheet held at arm's length and a slide read across a room are not the same
    piece of work. Passing `--min-pt` and its siblings still overrides both, and the line the check
    prints says which set it used, so a surprising result is traceable rather than mysterious.
    """
    deck = Presentation(str(deck_path))
    breite, hoehe = inches(deck.slide_width), inches(deck.slide_height)
    druck = _ist_druckstueck(breite, hoehe)
    vorgabe = ((PRINT_MIN_MARGIN, PRINT_MIN_PT, PRINT_MIN_PAD) if druck
               else (LAYOUT_MIN_MARGIN, LAYOUT_MIN_PT, LAYOUT_MIN_PAD))
    min_margin = vorgabe[0] if min_margin is None else min_margin
    min_pt = vorgabe[1] if min_pt is None else min_pt
    min_pad = vorgabe[2] if min_pad is None else min_pad
    slides = ([(slide_no, deck.slides[slide_no - 1])] if slide_no
              else list(enumerate(deck.slides, start=1)))
    problems, geprueft = [], 0
    for index, slide in slides:
        # Every leaf, not only the ones with words. Running off the page and sitting under
        # the margin are questions about a shape: a card, a logo or a rule carries no
        # text, and all three were invisible here, so a filled card hanging over the
        # edge passed clean.
        for shape, _rect in _leaf_rects(slide.shapes, mit_text=False):
            if shape.left is None or shape.width is None:
                continue
            geprueft += 1
            links, oben = inches(shape.left), inches(shape.top)
            rechts, unten = links + inches(shape.width), oben + inches(shape.height)
            if links < -0.02 or oben < -0.02 or rechts > breite + 0.02 or unten > hoehe + 0.02:
                problems.append(
                    "slide %d: '%s' runs off the slide (%.2f/%.2f to %.2f/%.2f on %.2f x %.2f)"
                    % (index, shape.name, links, oben, rechts, unten, breite, hoehe))
            elif min(links, oben, breite - rechts, hoehe - unten) < min_margin - 0.02:
                problems.append(
                    "slide %d: '%s' sits %.2f inch from the slide edge, under the %.2f floor"
                    % (index, shape.name, min(links, oben, breite - rechts, hoehe - unten),
                       min_margin))
            if not shape.has_text_frame or not text_of(shape):
                continue
            for groesse in sizes_in(shape):
                if groesse < min_pt - 0.01:
                    problems.append(
                        "slide %d: '%s' carries %g pt text, under the %g pt floor: '%s'"
                        % (index, shape.name, groesse, min_pt, text_of(shape)[:24]))
                    break
            gefuellt = False
            try:
                gefuellt = shape.fill.type == 1
            except Exception:  # noqa: BLE001 -- a shape without a fill type is simply not filled
                pass
            marker = inches(shape.width) < 1.0 and len(text_of(shape)) <= 4
            rand = inches(shape.text_frame.margin_left or 0)
            if gefuellt and not marker and rand < min_pad - 0.005:
                problems.append(
                    "slide %d: '%s' sets its text %.2f inch from its own filled edge, under the "
                    "%.2f floor" % (index, shape.name, rand, min_pad))
    for problem in problems:
        print(f"FAIL: {problem}")
    art = "print" if druck else "screen"
    print(f"{geprueft} shape(s) checked on {len(slides)} slide(s), {len(problems)} layout fault(s) "
          f"[{breite:.2f} x {hoehe:.2f}in, {art} floors: margin under {min_margin:g}in, "
          f"type under {min_pt:g}pt, padding under {min_pad:g}in]"
          + _nichts_geprueft(geprueft))
    return 1 if problems else 0


def _resolve_family(name: str | None, theme: dict | None) -> str | None:
    """The real typeface behind what a run names.

    A run in a deck built from a template does not carry "Montserrat", it carries `+mn-lt`, the
    theme's own placeholder for the body face. Passed on unchanged, no font file matches it, so
    every measurement returned None and every check that depends on one went quiet: `align-check`
    reported nothing on a built deck for the same reason `overflow-check` measured nothing. Found
    2026-08-26 while chasing the second of those.
    """
    if not name:
        return (theme or {}).get("minor")
    rolle = _THEME_PLACEHOLDER.get(name.strip().lower())
    if rolle:
        return (theme or {}).get(rolle)
    return name


def _theme_palette(deck_path: Path) -> set[str]:
    """Every colour the deck's first theme declares, uppercase hex without the hash."""
    import zipfile
    raus: set[str] = set()
    try:
        with zipfile.ZipFile(deck_path) as z:
            themes = sorted(n for n in z.namelist() if re.match(r"ppt/theme/theme\d+\.xml$", n))
            if not themes:
                return raus
            xml = z.read(themes[0]).decode("utf-8", "ignore")
    except (OSError, KeyError, zipfile.BadZipFile):
        return raus
    schema = re.search(r"<a:clrScheme.*?</a:clrScheme>", xml, re.S)
    if schema:
        for wert in re.findall(r'(?:srgbClr val|lastClr)="([0-9A-Fa-f]{6})"', schema.group(0)):
            raus.add(wert.upper())
    return raus


def _hard_colours(deck_path: Path) -> dict[str, int]:
    """Every colour the slides name outright, and how often. The theme is not consulted."""
    import zipfile
    gezaehlt: dict[str, int] = {}
    try:
        with zipfile.ZipFile(deck_path) as z:
            for name in z.namelist():
                if not re.match(r"ppt/slides/slide\d+\.xml$", name):
                    continue
                for wert in re.findall(r'srgbClr val="([0-9A-Fa-f]{6})"',
                                       z.read(name).decode("utf-8", "ignore")):
                    gezaehlt[wert.upper()] = gezaehlt.get(wert.upper(), 0) + 1
    except (OSError, KeyError, zipfile.BadZipFile):
        return {}
    return gezaehlt


# What a shape's look is made of, apart from where it sits. Deliberately not a list of properties
# this code understands: `a:xfrm` is position and stays with the migrated shape, everything else is
# the brand's and comes across whole. A property nobody has thought of yet travels with the rest.
_LOOK_ONLY = (qn("a:xfrm"),)


def _shape_look(element):
    """The look half of a shape: its `spPr` children except position, plus its `p:style`.

    Copying property by property is how this went wrong: the corner radius was handled, and then a
    shadow was missing, and after the shadow a line width. Every one of those is a rule somebody has
    to think of first, and the brand out there always has one more. Taking the whole thing instead
    means the list never has to be complete.
    """
    spPr = element.find(qn("p:spPr"))
    if spPr is None:
        return None
    look = [copy.deepcopy(k) for k in spPr if k.tag not in _LOOK_ONLY]
    style = element.find(qn("p:style"))
    return look, (copy.deepcopy(style) if style is not None else None)


def _brand_examples(deck_path: Path) -> dict:
    """One real example per kind of shape, taken from the brand's own slides.

    The kind is coarse on purpose: a filled box, an unfilled text box, a line. Finer than that and
    the match starts guessing what a shape means; coarser and a box gets a line's look.
    """
    beispiele: dict[str, list] = {}
    try:
        deck = Presentation(str(deck_path))
    except Exception:  # noqa: BLE001 -- an unreadable target is reported by the caller
        return {}
    for folie in deck.slides:
        for shape, _rect in _every_leaf(folie.shapes):
            el = shape._element
            spPr = el.find(qn("p:spPr"))
            if spPr is None or el.tag == qn("p:pic"):
                continue
            geom = spPr.find(qn("a:prstGeom"))
            if geom is None:
                continue
            gefuellt = spPr.find(qn("a:solidFill")) is not None or spPr.find(qn("a:gradFill")) is not None
            art = "line" if geom.get("prst") in ("line", "straightConnector1") else \
                  ("box" if gefuellt else "text")
            beispiele.setdefault(art, []).append(el)
    # The most common look of each kind wins, not the first one seen: a deck has one-off shapes and
    # the brand is what it does over and over.
    gewaehlt = {}
    for art, elemente in beispiele.items():
        zaehler: dict[str, list] = {}
        for el in elemente:
            spPr = el.find(qn("p:spPr"))
            from lxml import etree as _et
            schluessel = b"".join(_et.tostring(k) for k in spPr if k.tag not in _LOOK_ONLY).decode(
                "utf-8", "ignore")
            zaehler.setdefault(schluessel, []).append(el)
        haeufigste = max(zaehler.values(), key=len)
        if len(haeufigste) >= 3:
            gewaehlt[art] = (haeufigste[0], len(haeufigste), len(elemente))
    return gewaehlt


def _wireframe_furniture(slide) -> list[str]:
    """Shapes that carry no text and are too thin or too small to be content.

    A wireframe is a sketch of where things go, and it draws its own aids to show that: rules beside
    a block, ticks, spacers, boxes standing in for a picture. Those are the sketch's language, not
    the piece's, and `migrate` brings the whole arrangement across including them. Reported from the
    field nine 0.04 inch bars ended up as rules in front of every text block, in a brand
    whose own 74 vertical rules only ever sit between an icon and its text column. The shape was
    formally right and in the wrong role, which no geometry check can see.

    Listed, never removed: a thin bar can be exactly right where the brand uses one, and only a
    person can say which it is.
    """
    raus = []
    for shape, (_l, _t, breite, hoehe) in _every_leaf(slide.shapes):
        if shape.has_text_frame and text_of(shape):
            continue
        if shape._element.tag == qn("p:pic") or not breite or not hoehe:
            continue
        kurz, lang = sorted((inches(breite), inches(hoehe)))
        # The short side is what tells a rule from a band. A band 11.5 inches wide and 0.64 high has
        # a ratio of 18 to 1 and is a deliberate element; a rule 0.04 wide is furniture whatever its
        # length. Zero-sized leftovers (think-cell markers and the like) are not shapes at all.
        if kurz <= 0.01 or lang <= 0.01:
            continue
        if kurz < 0.15:
            raus.append(f"{shape.name} ({inches(breite):.2f} x {inches(hoehe):.2f} inch)")
    return raus


def _apply_brand_look(slide, beispiele: dict) -> tuple[int, list[str]]:
    """Give every migrated shape the whole look of its kind in the brand. Returns count and detail."""
    geaendert, gezaehlt = 0, {}
    for shape, _rect in _every_leaf(slide.shapes):
        el = shape._element
        spPr = el.find(qn("p:spPr"))
        if spPr is None or el.tag == qn("p:pic"):
            continue
        geom = spPr.find(qn("a:prstGeom"))
        if geom is None:
            continue
        gefuellt = spPr.find(qn("a:solidFill")) is not None or spPr.find(qn("a:gradFill")) is not None
        art = "line" if geom.get("prst") in ("line", "straightConnector1") else \
              ("box" if gefuellt else "text")
        if art not in beispiele:
            continue
        vorbild = _shape_look(beispiele[art][0])
        if vorbild is None:
            continue
        look, style = vorbild
        for kind in list(spPr):
            if kind.tag not in _LOOK_ONLY:
                spPr.remove(kind)
        for kind in look:
            spPr.append(copy.deepcopy(kind))
        alter_style = el.find(qn("p:style"))
        if alter_style is not None:
            el.remove(alter_style)
        if style is not None:
            spPr.addnext(copy.deepcopy(style))
        geaendert += 1
        gezaehlt[art] = gezaehlt.get(art, 0) + 1
    detail = [f"{n}x {art}" for art, n in sorted(gezaehlt.items())]
    return geaendert, detail


def _theme_carries_the_brand(deck_path: Path) -> str | None:
    """A warning where the target's theme palette is not the brand its own slides paint in.

    `migrate` maps onto theme roles, and the whole promise of that is "the target's theme colours
    all of it". Where the target deck carries a stock theme -- an untouched Office palette under a
    template name -- that promise turns into its opposite: the migrated pattern comes out in Office
    magenta while every slide around it is painted in a brand colour the theme has never heard of.
    Seen in practice: 1186 hard #00005A against a theme that does not know it.

    Read from the file, not from the object, and only reported. Which colour is the brand is not a
    thing a script may decide.
    """
    palette = _theme_palette(deck_path)
    hart = _hard_colours(deck_path)
    if not palette or not hart:
        return None
    haeufig = [(w, z) for w, z in sorted(hart.items(), key=lambda kv: -kv[1])
               if z >= 20 and w not in palette and w not in ("FFFFFF", "000000")]
    if not haeufig:
        return None
    genannt = ", ".join(f"#{w} ({z}x)" for w, z in haeufig[:3])
    return (f"this deck's own slides paint in {genannt}, and its theme palette holds none of them. "
            f"Mapping onto theme roles therefore does not adopt this brand, it replaces it with "
            f"whatever the theme carries. Check the result's colours, or give the deck a theme "
            f"that is the brand.")


def _theme_fonts(deck_path: Path) -> dict:
    """The major and minor latin typefaces this deck's theme declares."""
    import zipfile
    try:
        with zipfile.ZipFile(deck_path) as z:
            themes = sorted(n for n in z.namelist() if re.match(r"ppt/theme/theme\d+\.xml$", n))
            if not themes:
                return {}
            xml = z.read(themes[0]).decode("utf-8", "ignore")
    except (OSError, KeyError, zipfile.BadZipFile):
        return {}
    out = {}
    for rolle, tag in (("major", "majorFont"), ("minor", "minorFont")):
        treffer = re.search(r"<a:%s>.*?<a:latin[^>]*typeface=\"([^\"]+)\"" % tag, xml, re.S)
        if treffer and treffer.group(1):
            out[rolle] = treffer.group(1)
    return out


def _theme_font(deck_path: Path) -> str | None:
    """The body typeface this deck's own theme names.

    A run that does not name a font is not fontless: it inherits, and what it inherits is what
    PowerPoint reads out of the theme. Without this the measurement gave up on every built deck,
    where the runs carry no explicit family, which is exactly the file that needs checking. This
    reads the deck's own declaration, so it stays measured rather than assumed.
    """
    import zipfile
    try:
        with zipfile.ZipFile(deck_path) as z:
            themes = sorted(n for n in z.namelist() if re.match(r"ppt/theme/theme\d+\.xml$", n))
            if not themes:
                return None
            xml = z.read(themes[0]).decode("utf-8", "ignore")
    except (OSError, KeyError, zipfile.BadZipFile):
        return None
    treffer = re.search(r'<a:minorFont>.*?<a:latin[^>]*typeface="([^"]+)"', xml, re.S)
    return treffer.group(1) if treffer and treffer.group(1) else None


def _frame_height(rahmen, breite_emu, groessen, fallback_family: dict | None = None) -> float | None:
    """Inches of height a text frame really needs in a given width, its own margins included.

    Counting lines and multiplying by one size was the earlier arithmetic, and it is wrong wherever
    a box holds more than one size, which is most cards: a heading at 16 pt over a paragraph at
    12 pt got all six lines charged at 16 and then compared against a box that had room for that,
    while the render showed the text running out of its area. Paragraph spacing was not in it at
    all. Seen in practice, on three cards a check called clean.

    Each paragraph is measured at its own size, with its own line spacing and the space before and
    after it. Where a paragraph cannot be measured the whole frame returns None, never a number
    that is too small.
    """
    hoehe = inches(rahmen.margin_top if rahmen.margin_top is not None else 45720) \
        + inches(rahmen.margin_bottom if rahmen.margin_bottom is not None else 45720)
    zeilen = _frame_lines(rahmen, breite_emu, groessen, fallback_family, hoehe_aus=hoehe)
    return zeilen


def _frame_lines(rahmen, breite_emu, groessen, fallback_family: dict | None = None,
                 hoehe_aus: float | None = None) -> float | None:
    """Lines a text frame really takes in a given width. Shared by shapes and table cells, because
    a cell overflows exactly the way a box does and the arithmetic must not differ between them."""
    links = inches(rahmen.margin_left if rahmen.margin_left is not None else 91440)
    rechts = inches(rahmen.margin_right if rahmen.margin_right is not None else 91440)
    breite_in = inches(breite_emu) - links - rechts
    if breite_in <= 0:
        return None
    ersatz_groesse = max(groessen or [0.0])
    zeilen = 0.0
    hoehe = hoehe_aus or 0.0
    for paragraph in rahmen.paragraphs:
        text = "".join(r.text for r in paragraph.runs)
        abstand = _paragraph_spacing(paragraph)
        if not text.strip():
            zeilen += 1
            hoehe += (ersatz_groesse or 12.0) * 1.2 / 72.0 + abstand[0] + abstand[1]
            continue
        size = max((r.font.size.pt for r in paragraph.runs if r.font.size), default=0.0) or ersatz_groesse
        family = _resolve_family(
            next((r.font.name for r in paragraph.runs if r.font.name), None), fallback_family)
        bold = any(r.font.bold for r in paragraph.runs)
        if not size:
            return None
        real = _real_ink_bbox(text, family, size, bold)
        if real is None:
            return None
        # A line breaks between words, and a word that fits stays whole: without the cap below, a
        # single wide glyph, a "2" at 54 pt in a narrow box, counted as two lines and was reported
        # as an overflow it cannot have, since one character cannot wrap.
        #
        # But a *word* wider than its own column does wrap, in the middle, and renderers do it. The
        # cap on the word count then hid the case entirely: "Systemaufnahme" at 20 pt bold measures
        # 2.28 inch in a 2.05 inch column, was counted as one line, the box was sized for one, and
        # the second line landed on whatever sat below with every check reporting clean. German
        # compounds in narrow columns are the normal case for a flyer, not an edge case.
        #
        # So the cap is the word count plus the breaks that happen inside over-wide words. Measured
        # per word only where the totals disagree, because that measurement is the expensive part
        # and it is pointless where nothing is too wide.
        woerter = text.split()
        gebrochen = math.ceil(real[1] / breite_in)
        deckel = float(len(woerter))
        if gebrochen > len(woerter):
            for wort in woerter:
                if len(wort) < 2:
                    continue      # one character cannot break; it paints past the edge instead
                masse = _real_ink_bbox(wort, family, size, bold)
                if masse and masse[1] > breite_in:
                    deckel += math.ceil(masse[1] / breite_in) - 1
        dieser = max(1.0, min(float(gebrochen), deckel))
        zeilen += dieser
        hoehe += dieser * size * _line_spacing(paragraph) / 72.0 + abstand[0] + abstand[1]
    if hoehe_aus is not None:
        return hoehe or None
    return zeilen or None


def _line_spacing(paragraph) -> float:
    """The paragraph's own line spacing as a factor, 1.2 where it names none."""
    pPr = paragraph._pPr
    if pPr is None:
        return 1.2
    ln = pPr.find(qn("a:lnSpc"))
    if ln is None:
        return 1.2
    pct = ln.find(qn("a:spcPct"))
    if pct is not None and pct.get("val"):
        return int(pct.get("val")) / 100000.0
    pts = ln.find(qn("a:spcPts"))
    if pts is not None and pts.get("val"):
        groessen = [r.font.size.pt for r in paragraph.runs if r.font.size]
        basis = max(groessen) if groessen else 12.0
        return (int(pts.get("val")) / 100.0) / basis if basis else 1.2
    return 1.2


def _paragraph_spacing(paragraph) -> tuple[float, float]:
    """(before, after) in inches. Points and percent both occur; percent is of the type size."""
    raus = []
    pPr = paragraph._pPr
    groessen = [r.font.size.pt for r in paragraph.runs if r.font.size]
    basis = max(groessen) if groessen else 12.0
    for tag in ("a:spcBef", "a:spcAft"):
        wert = 0.0
        knoten = pPr.find(qn(tag)) if pPr is not None else None
        if knoten is not None:
            pts = knoten.find(qn("a:spcPts"))
            pct = knoten.find(qn("a:spcPct"))
            if pts is not None and pts.get("val"):
                wert = int(pts.get("val")) / 100.0 / 72.0
            elif pct is not None and pct.get("val"):
                wert = basis * (int(pct.get("val")) / 100000.0) / 72.0
        raus.append(wert)
    return raus[0], raus[1]


def _lines_needed(shape, fallback_family: dict | None = None) -> float | None:
    """How many lines this shape's text really takes in its own box width, or None if it cannot
    be measured on this machine.

    Measured, not counted: a character count says "Plattform- & Virtualisierungsstrategien" fits a
    7.09 inch box at 26 pt, and the render says it wraps. The width comes from the same font file
    the deck asks for, through `_real_ink_bbox`, which is what `align_check` already trusts. Where
    the font cannot be resolved this returns None and the shape is left out, never assumed to fit.
    """
    # The insets, not the box. PowerPoint lays the text out inside the frame's margins, 0.1 inch
    # left and right by default, and the line breaks against that width. Measured on the reported
    # case: the title needs 6.96 inch, the box is 7.09 wide and looks roomy, and the text area is
    # 6.89, so it wrapped. Ignoring the insets makes the check agree with the geometry and disagree
    # with the render, which is the failure it exists to catch.
    return _frame_lines(shape.text_frame, shape.width, sizes_in(shape), fallback_family)


# How much of a filled shape its own text has to occupy before it stops looking like a mistake.
# Seventy per cent is where a card reads as set rather than as half-filled: measured on a real
# piece, a results card of 4.15 cm carried text needing 2.93 cm and left a visible hole under the
# second row, and every check passed it, because nothing overflowed, nothing overlapped and nothing
# sat outside the page. All of them ask whether what is there is wrong. None asks whether something
# is missing.
FILL_MIN = 0.70


def fill_check(deck_path: Path, slide_no: int | None = None, minimum: float = None) -> int:
    """A filled shape using far less of itself than it has, which is either a fault or a pause.

    The mirror image of `overflow_check`, and the half nobody built. That one asks whether the text
    is too big for the box; this asks whether the box is too big for the text. The second is harder
    to see, because nothing about it is broken: the file is valid, the render is clean, and it
    simply looks unfinished.

    **The text of a card is rarely inside the card.** A first version asked each filled shape about
    its own text frame and found nothing at all across nine finished flyers: the cards were filled
    auto shapes with no text, and the text sat in separate boxes laid over them, which themselves
    carry no fill. Each card fell through both meshes. That layout is not exotic either, it is what
    comes out of following the rules in `design.md` about setting colour explicitly. So the
    question is asked geometrically instead: what stands inside this shape, whether or not it is
    part of it.

    Only filled shapes are looked at. The same empty space around a free-standing line is the white
    space the piece is made of, not a hole in something. A shape whose text cannot be measured is
    left out rather than called clean.
    """
    grenze = FILL_MIN if minimum is None else minimum
    deck = Presentation(str(deck_path))
    slides = ([(slide_no, deck.slides[slide_no - 1])] if slide_no
              else list(enumerate(deck.slides, start=1)))
    theme = _theme_fonts(deck_path)
    problems, geprueft = [], 0
    rand = round(0.02 * EMU_PER_INCH)   # a hairline of slack, so a box on the edge still counts

    for index, slide in slides:
        # Every leaf, not only the ones carrying words: the card is the shape with no text in it.
        blaetter = list(_leaf_rects(slide.shapes, mit_text=False))
        for shape, (links, oben, breite, hoehe_emu) in blaetter:
            try:
                gefuellt = shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE or shape.fill.type == 1
                if shape.fill.type == 5:      # background fill: deliberately not a card
                    gefuellt = False
            except (AttributeError, TypeError, ValueError):
                continue
            if not gefuellt:
                continue
            hoch = inches(hoehe_emu)
            if hoch < 0.4:                    # a strip or a label, not something with room in it
                continue
            # The inset is not emptiness. A card that a kit gives half a centimetre of padding top
            # and bottom has that much less room by design, and measuring the text against the
            # outer height reports the padding as missing content. On a 3.4 cm card two 0.5 cm
            # insets are a third of it, so every small card came out at about two thirds full and
            # the check fired systematically on exactly the ones that were correct.
            innen = 0.0
            if shape.has_text_frame:
                rahmen = shape.text_frame
                oben_in = rahmen.margin_top if rahmen.margin_top is not None else 45720
                unten_in = rahmen.margin_bottom if rahmen.margin_bottom is not None else 45720
                innen = inches(oben_in) + inches(unten_in)
            nutzbar = max(hoch - innen, 0.05)

            eigener = _frame_height(shape.text_frame, breite, sizes_in(shape), theme) \
                if shape.has_text_frame and text_of(shape) else None
            if eigener is not None:
                belegt = eigener
            else:
                # What stands inside this shape's box. Measured from the top of the highest thing
                # in it to the bottom of what the lowest thing's text actually needs, because a box
                # is often taller than its own text and the gap between two of them is part of the
                # card rather than empty space at its end.
                oberkante, unterkante = None, None
                for andere, (l2, o2, b2, h2) in blaetter:
                    if andere is shape or not andere.has_text_frame or not text_of(andere):
                        continue
                    if not (l2 >= links - rand and o2 >= oben - rand
                            and l2 + b2 <= links + breite + rand
                            and o2 + h2 <= oben + hoehe_emu + rand):
                        continue
                    noetig = _frame_height(andere.text_frame, b2, sizes_in(andere), theme)
                    if noetig is None:
                        continue
                    o_in = inches(o2)
                    oberkante = o_in if oberkante is None else min(oberkante, o_in)
                    unterkante = max(unterkante or 0.0, o_in + noetig)
                if oberkante is None:
                    continue
                belegt = (unterkante - oberkante) + (oberkante - inches(oben))
            geprueft += 1
            anteil = belegt / nutzbar if nutzbar else 1.0
            if anteil < grenze:
                was = text_of(shape)[:24] if (shape.has_text_frame and text_of(shape)) else shape.name
                problems.append(
                    "slide %d: '%s' is filled to %.0f%% of the room it has (%.2f of %.2f inch, "
                    "%.2f of that inset). Either content is missing or the breathing room is "
                    "meant; say which."
                    % (index, was, anteil * 100, belegt, nutzbar, innen))
    for problem in problems:
        print(f"FAIL: {problem}")
    print(f"{geprueft} filled shape(s) checked on {len(slides)} slide(s), {len(problems)} "
          f"under {grenze:.0%} full" + _nichts_geprueft(geprueft))
    return 1 if problems else 0


def overflow_check(deck_path: Path, slide_no: int | None, baseline: Path | None = None,
                   als_liste: bool = False):
    """Text that needs more height than its own box has.

    `overlap-check` compares the boxes two shapes declare, so it is blind to the case where one
    shape's text wraps past its own box and lands on the shape below: the geometry still says they
    do not touch, and the render says otherwise. Found on a real deck, where a 26 pt
    title wrapped to a second line and sat on the claim under it, letters on letters, while
    overlap-check reported 1134 pairs and no overlap. Only looking at the render caught it.

    Two things this deliberately does not do. It does not use `capacity()`, the character count
    `build` uses to decide whether text fits a slot: measured against a finished deck that produced
    31 findings on a reference deck its owner had corrected by hand, and missed the one case that
    was actually broken. And it stays quiet where the text shrinks to fit, since nothing can
    overflow there. A shape whose font cannot be measured is left out, never called clean.
    """
    deck = Presentation(str(deck_path))
    slides = [(slide_no, deck.slides[slide_no - 1])] if slide_no else list(enumerate(deck.slides, start=1))
    theme = _theme_fonts(deck_path)
    problems, checked, unmeasured = [], 0, 0
    nicht_messbar: list[str] = []
    for index, slide in slides:
        # The real rectangle, not the declared one: a shape inside a group states its size in the
        # group's own units, and measuring that against its text compares two different scales.
        for shape, (_l, _t, breite, hoehe_emu) in _leaf_rects(slide.shapes):
            if not shape.has_text_frame or not text_of(shape):
                continue
            if shape.text_frame.auto_size == MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE:
                continue
            zeilen = _frame_lines(shape.text_frame, breite, sizes_in(shape), theme)
            noetig = _frame_height(shape.text_frame, breite, sizes_in(shape), theme)
            if zeilen is None or noetig is None:
                unmeasured += 1
                nicht_messbar.append(f"slide {index}: '{text_of(shape)[:24]}'")
                continue
            checked += 1
            size = max(sizes_in(shape) or [12.0])
            hoch = inches(hoehe_emu)
            # Only a wrap counts. A single line in a box a tenth of an inch shorter than its own
            # line height is a design choice and paints fine; reporting it buried the one finding
            # that mattered under a symbol in a tight box. The damage this check exists for is the
            # second line, which lands on whatever sits below.
            if zeilen > 1 and noetig > hoch + 0.02:
                problems.append(
                    "slide %d: '%s' needs about %.2f inch for %d line(s) at up to %.0f pt, its box "
                    "is %.2f inch high"
                    % (index, text_of(shape)[:24], noetig, int(zeilen), size, hoch))
        # A battlecard is mostly a table, and a cell overflows exactly the way a box does: the
        # fourth line of a pitch lands on the row below and takes its headings with it. Nothing saw
        # that, because every check here walks shapes and a cell is not one, so headings can sit
        # hidden in a render on a file that was already approved.
        # A battlecard is mostly a table, and nothing here saw inside one: every check walks shapes,
        # and a cell is not a shape. Seen in practice, on a file the user had
        # already approved, where a pitch ran to a fourth line and covered the two headings below.
        #
        # Measured with `cell_capacity`, not with the line arithmetic above. Two attempts at
        # measuring a cell by height were thrown away first: a row's stated height is a minimum,
        # PowerPoint grows the row and pushes the rest down, so height says almost nothing. Both
        # attempts reported an overflow on a battlecard that renders correctly, and neither told the
        # reported case apart from the clean one. `cell_capacity` is calibrated against real decks
        # and is what `build` already refuses on; the reported cell held 568 characters in a slot
        # measured at 372, which this catches and height did not.
        for shape in slide.shapes:
            if not getattr(shape, "has_table", False):
                continue
            table = shape.table
            for r in range(len(table.rows)):
                for c in range(len(table.columns)):
                    zelle = table.cell(r, c)
                    if zelle.is_spanned or not zelle.text_frame.text.strip():
                        continue
                    breite = table.columns[c].width or 0
                    hoehe = table.rows[r].height or 0
                    if not breite or not hoehe:
                        continue
                    raum = cell_capacity(zelle, breite, hoehe,
                                         _cell_family(zelle, theme))
                    if raum <= 0:
                        continue
                    checked += 1
                    text = zelle.text_frame.text.strip()
                    if len(text) > raum:
                        problems.append(
                            "slide %d: cell r%dc%d holds about %d characters and carries %d, so it "
                            "grows the row and covers what sits below: '%s'"
                            % (index, r + 1, c + 1, raum, len(text), text[:40]))
    if als_liste:
        return problems
    problems, geerbt = _nur_neue(
        problems, baseline, lambda p: overflow_check(p, None, als_liste=True))
    for problem in problems:
        print(f"FAIL: {problem}")
    rest = f", {unmeasured} not measurable" if unmeasured else ""
    # Named, not just counted. What cannot be measured is passed over silently otherwise, and a
    # silent pass reads exactly like a clean one.
    for eintrag in nicht_messbar[:12]:
        print(f"  not measurable: {eintrag}")
    if len(nicht_messbar) > 12:
        print(f"  ... and {len(nicht_messbar) - 12} more")
    if baseline:
        rest += f", {geerbt} already in {baseline.name}"
    print(f"{checked} text box(es) and cell(s) checked on {len(slides)} slide(s), "
          f"{len(problems)} overflow(s){rest}")
    return 1 if problems else 0


def slots(deck_path: Path, slide_no: int | None) -> int:
    """The fillable places in a deck that already exists, by name, with what stands there now.

    `show` does this for a harvested library. This does it for any file, which is the case that was
    missing: changing the wording of a finished deck needed either a full rebuild from a library or
    a purpose-written python-pptx script, and the second is what kept happening. The names printed
    here are what `fill` takes.
    """
    deck = Presentation(str(deck_path))
    paare = [(slide_no, deck.slides[slide_no - 1])] if slide_no else list(enumerate(deck.slides, start=1))
    gesamt = 0
    for index, slide in paare:
        gefunden = slot_names(slide)
        print(f"Slide {index}: {len(gefunden)} slot(s)")
        for name in sorted(gefunden):
            eintrag = gefunden[name]
            text = " ".join(str(eintrag["text"]).split())
            print(f"   {name:<22} holds ~{eintrag['capacity']:>4} chars   now {len(eintrag['text']):>4}"
                  f"   {text[:48]}")
        gesamt += len(gefunden)
        print()
    print(f"{gesamt} slot(s) on {len(paare)} slide(s)")
    return 0


def fill(deck_path: Path, texts_path: Path, out: Path | None, strict: bool) -> int:
    """Change the wording of a deck that already exists, in place or into a copy.

    The whole mechanic for this was already here, in `fill_slots`: address a place by name, refuse
    text that does not fit rather than shrink it. What was missing was a way in without a library,
    so every wording change on a finished file became a one-off script. Written after
    that had happened twice in one bundle.

    The texts file is `{"1": {"slot": "new text"}}` by slide number, or a flat `{"slot": "text"}`
    when the deck has one slide. A list of strings fills several paragraphs.
    """
    daten = json.loads(texts_path.read_text(encoding="utf-8"))
    deck = Presentation(str(deck_path))
    if daten and not any(str(k).isdigit() for k in daten):
        if len(deck.slides) != 1:
            print("slide-library: this deck has more than one slide, so the texts file has to name "
                  "them: {\"1\": {...}, \"2\": {...}}", file=sys.stderr)
            return 2
        daten = {"1": daten}
    problems = []
    angefragt = 0
    for schluessel, texte in daten.items():
        try:
            nummer = int(schluessel)
            slide = deck.slides[nummer - 1]
        except (ValueError, IndexError):
            print(f"slide-library: there is no slide {schluessel} in this deck", file=sys.stderr)
            return 2
        angefragt += len(texte)
        problems += [f"slide {nummer}: {p}" for p in fill_slots(slide, texte, {}, strict)]
    ziel = out or deck_path
    if not guarded_save(deck, ziel):
        return 1
    rest = f", {len(problems)} reported" if problems else ""
    print(f"{angefragt} slot(s) asked for, written to {ziel}{rest}")
    lose = dangling_refs(ziel)
    for eintrag in lose:
        print(f"FAIL: {eintrag}", file=sys.stderr)
    if lose:
        return 1
    for problem in problems:
        print(f"FAIL: {problem}", file=sys.stderr)
    if problems:
        print(f"{len(problems)} open. Shorten the text, or pass --loose to write it anyway.",
              file=sys.stderr)
        return 1
    return 0


def main(argv: list[str]) -> int:
    ap = argparse.ArgumentParser(description="Harvest a brand's own slides, then fill copies of them.")
    sub = ap.add_subparsers(dest="command", required=True)

    ph = sub.add_parser("harvest", help="read a deck and write the library index")
    ph.add_argument("deck", type=Path)
    ph.add_argument("--into", type=Path, required=True)

    pb = sub.add_parser("build", help="clone slides from the library and fill them")
    pb.add_argument("plan", type=Path)
    pb.add_argument("out", type=Path)
    pb.add_argument("--library", type=Path, required=True)
    pb.add_argument("--loose", action="store_true",
                    help="write text that exceeds a slot's capacity anyway, still reported")

    ps = sub.add_parser("show", help="print the library so a slide can be chosen")
    ps.add_argument("library", type=Path)
    ps.add_argument("--slide", type=int)

    pc = sub.add_parser("check", help="print the library and record that it was looked at for a task")
    pc.add_argument("library", type=Path)
    pc.add_argument("--task", required=True,
                     help="the workbench/<slug>/ bundle this deliverable belongs to")
    pc.add_argument("--space", type=Path, default=Path.cwd(),
                     help="space root the task's zanmai/temp/ lives under (default: cwd)")
    pc.add_argument("--shape", help="the pattern chosen for this piece, recorded for the bundle")
    pc.add_argument("--why", help="why this content has that shape")

    pn = sub.add_parser("nudge", help="move one shape by a delta, all four xfrm values set explicitly")
    pn.add_argument("deck", type=Path)
    pn.add_argument("--shape", required=True, help="shape name, or full slash path when the name repeats")
    pn.add_argument("--dx", type=float, default=0.0, help="inches, positive is right")
    pn.add_argument("--dy", type=float, default=0.0, help="inches, positive is down")
    pn.add_argument("--slide", type=int, help="1-based; default searches every slide, errors if ambiguous")
    pn.add_argument("--into", type=Path, required=True)

    psl = sub.add_parser("slots", help="the fillable places in a deck that already exists")
    psl.add_argument("deck", type=Path)
    psl.add_argument("--slide", type=int)

    pfi = sub.add_parser("fill", help="change the wording of a deck that already exists")
    pfi.add_argument("deck", type=Path)
    pfi.add_argument("--texts", type=Path, required=True,
                     help='JSON: {"1": {"slot": "new text"}}, or flat for a one-slide deck')
    pfi.add_argument("--out", type=Path, help="write a copy instead of changing the file itself")
    pfi.add_argument("--loose", action="store_true",
                     help="write text that exceeds a slot's capacity anyway, still reported")

    pr = sub.add_parser("refs-check", help="relationship ids a slide points at that its rels do not carry")
    pr.add_argument("deck", type=Path)

    pfc = sub.add_parser("fill-check", help="a filled shape using far less of its own height than it has: content missing, or breathing room meant")
    pfc.add_argument("deck", type=Path)
    pfc.add_argument("--slide", type=int, help="1-based; default checks every slide")
    pfc.add_argument("--min-fill", type=float, dest="min_fill",
                     help=f"fraction of its own height a filled shape should use; default {FILL_MIN}")

    pmc = sub.add_parser("media-check", help="every picture the deck points at, checked for being a real file")
    pmc.add_argument("deck", type=Path)
    pmc.add_argument("--verbose", action="store_true", help="name every picture, not only the broken ones")

    psc = sub.add_parser("schema-check", help="shapes PowerPoint will not draw, though a render shows them: "
                                             "a:spPr out of schema order, or a shape with no position")
    psc.add_argument("deck", type=Path)
    psc.add_argument("--slide", type=int, help="1-based; default checks every slide")

    pf = sub.add_parser("overflow-check", help="text that needs more room than its own box has")
    pf.add_argument("deck", type=Path)
    pf.add_argument("--slide", type=int)
    pf.add_argument("--baseline", type=Path,
                    help="the deck this was built from; findings it already has are not reported")

    po = sub.add_parser("overlap-check", help="pairwise text-over-text check, ink-aware for wrap=none")
    po.add_argument("deck", type=Path)
    po.add_argument("--slide", type=int, help="1-based; default checks every slide")
    po.add_argument("--baseline", type=Path,
                    help="the deck this was built from; findings it already has are not reported")

    pa = sub.add_parser("align-check", help="ink-left alignment between text frames sharing a box-left edge")
    pa.add_argument("deck", type=Path)
    pa.add_argument("--slide", type=int, help="1-based; default checks every slide")

    pf = sub.add_parser("furniture-check", help="an element that repeats across pages sitting somewhere else on one of them")
    pf.add_argument("deck", type=Path)
    pf.add_argument("--tol", type=float, help=f"inch of drift allowed; default {FURNITURE_TOL}")
    pf.add_argument("--zone", type=float, help=f"fraction of page height at top and bottom that counts as furniture; default {FURNITURE_ZONE}")

    pr = sub.add_parser("render", help="a picture of every slide, headless, any platform")
    pr.add_argument("deck", type=Path)
    pr.add_argument("--into", type=Path, required=True, help="directory for the pictures")
    pr.add_argument("--font-dir", type=Path, action="append", dest="font_dirs",
                    help="a folder of typefaces this piece is set in, repeatable. Use it where the "
                         "face lives with the brand rather than being installed on the machine.")
    pr.add_argument("--dpi", type=int, default=110)

    pst = sub.add_parser("structure-check", help="is the built slide still carrying the structure "
                                                "of the pattern it was built from")
    pst.add_argument("deck", type=Path)
    pst.add_argument("--slide", type=int, required=True)
    pst.add_argument("--against", required=True, help="<pattern deck>:<slide>")
    pst.add_argument("--intended", help="why the build deliberately differs from the pattern; "
                                       "a wireframe is a starting point, not a template")
    pst.add_argument("--tolerance", type=float, default=0.0,
                     help="how much less the build may hold; zero by default, because a proportion "
                          "says nothing about whether what is missing was load-bearing")

    pl2 = sub.add_parser("leftover-check", help="what the file carries beyond its content: "
                                              "comments, speaker notes, animations, authors")
    pl2.add_argument("deck", type=Path)

    psw = sub.add_parser("swap-image", help="replace a picture with another, whole picture "
                                           "(crop included), keeping the old one's placement")
    psw.add_argument("deck", type=Path)
    psw.add_argument("--shape", required=True, help="name or full path of the picture to replace")
    psw.add_argument("--from", dest="quelle", required=True, help="<deck>:<slide>:<shape>")
    psw.add_argument("--slide", type=int, help="1-based, where the target picture sits")
    psw.add_argument("--out", type=Path, required=True)

    psg = sub.add_parser("suggest", help="which patterns carry a content of this shape; a shortlist, "
                                        "not a decision")
    psg.add_argument("library", type=Path, nargs="?",
                     default=Path("zanmai/system/templates/wireframes"))
    psg.add_argument("--elements", type=int, help="how many things the content holds")
    psg.add_argument("--order", choices=["equal", "sequence", "one-over-many", "matrix", "pairs",
                                         "single"],
                     help="equal rank, an order, one over many, two axes, two sides, or one piece")
    psg.add_argument("--movement", choices=["yes", "no"],
                     help="does something move from A to B")

    pk = sub.add_parser("keep", help="put an approved slide into the brand's library, so the next "
                                     "one of its kind is a copy")
    pk.add_argument("deck", type=Path)
    pk.add_argument("--slide", type=int, required=True)
    pk.add_argument("--brand-dir", dest="brand_dir", type=Path, required=True,
                    help="the brand's folder, e.g. zanmai/design/<brand>")
    pk.add_argument("--as", dest="slug", required=True, help="short name for this shape of content")
    pk.add_argument("--source", help="where it came from, for the note beside it")

    px = sub.add_parser("extract", help="lift whole slides out of a deck, XML untouched, links "
                                       "into dropped slides cut")
    px.add_argument("deck", type=Path)
    px.add_argument("--slides", required=True,
                    help="1-based slide numbers, comma-separated, in the order they should end up")
    px.add_argument("--out", type=Path, required=True)

    pm = sub.add_parser("migrate", help="put one slide into another deck, adopting its master, "
                                       "theme and layout")
    pm.add_argument("deck", type=Path, help="the deck the slide comes from")
    pm.add_argument("--slide", type=int, required=True, help="1-based slide number in that deck")
    pm.add_argument("--into", type=Path, required=True, help="the deck whose brand it should adopt")
    pm.add_argument("--out", type=Path, required=True)
    pm.add_argument("--layout", help="layout name on the target's brand master")
    pm.add_argument("--scheme", help="colour scheme name that identifies the brand master")
    pm.add_argument("--title", help="text for the layout's title placeholder")
    pm.add_argument("--replace", action="store_true",
                    help="empty the target first; without it the slide is appended")
    pm.add_argument("--keep-shapes", dest="keep_shapes", action="store_true",
                    help="leave the source's look alone instead of taking the brand's")
    pm.add_argument("--brand-from", dest="brand_from", type=Path,
                    help="deck the brand is measured in, where the target is still empty or thin")

    pl = sub.add_parser("layout-check", help="shapes off the slide, under the margin, "
                                             "under the type floor, or touching their own fill")
    pl.add_argument("deck", type=Path)
    pl.add_argument("--slide", type=int, help="1-based; default checks every slide")
    pl.add_argument("--min-margin", type=float, help=f"inch from the slide edge (default {LAYOUT_MIN_MARGIN})")
    pl.add_argument("--min-pt", type=float, help=f"smallest type allowed (default {LAYOUT_MIN_PT})")
    pl.add_argument("--min-pad", type=float, help=f"inch inside a filled box (default {LAYOUT_MIN_PAD})")

    args = ap.parse_args(argv[1:])
    if args.command == "harvest":
        if not args.deck.is_file():
            print(f"slide-library: no deck at {args.deck}", file=sys.stderr)
            return 2
        return harvest(args.deck, args.into)
    if args.command == "build":
        return build(args.plan, args.out, args.library, strict=not args.loose)
    if args.command == "check":
        return check(args.library, args.task, args.space, args.shape, args.why)
    if args.command in ("nudge", "overlap-check", "align-check", "layout-check") and not args.deck.is_file():
        print(f"slide-library: no deck at {args.deck}", file=sys.stderr)
        return 2
    if args.command == "nudge":
        return nudge(args.deck, args.shape, args.dx, args.dy, args.slide, args.into)
    if args.command == "slots":
        return slots(args.deck, args.slide)
    if args.command == "fill":
        return fill(args.deck, args.texts, args.out, strict=not args.loose)
    if args.command == "fill-check":
        return fill_check(args.deck, args.slide, args.min_fill)
    if args.command == "media-check":
        return media_check(args.deck, args.verbose)
    if args.command == "refs-check":
        lose = dangling_refs(args.deck)
        for eintrag in lose:
            print(f"FAIL: {eintrag}")
        print(f"{len(lose)} reference(s) into nothing")
        return 1 if lose else 0
    if args.command in ("overflow-check", "overlap-check") and args.baseline \
            and not args.baseline.is_file():
        print(f"slide-library: no baseline deck at {args.baseline}", file=sys.stderr)
        return 2
    if args.command == "structure-check":
        if not args.deck.is_file():
            print(f"slide-library: no deck at {args.deck}", file=sys.stderr)
            return 2
        return structure_check(args.deck, args.slide, args.against, args.tolerance, args.intended)
    if args.command == "leftover-check":
        if not args.deck.is_file():
            print(f"slide-library: no deck at {args.deck}", file=sys.stderr)
            return 2
        return leftover_check(args.deck)
    if args.command == "swap-image":
        if not args.deck.is_file():
            print(f"slide-library: no deck at {args.deck}", file=sys.stderr)
            return 2
        return swap_image(args.deck, args.shape, args.quelle, args.slide, args.out)
    if args.command == "suggest":
        return suggest(args.library, args.elements, args.order, args.movement)
    if args.command == "keep":
        if not args.deck.is_file():
            print(f"slide-library: no deck at {args.deck}", file=sys.stderr)
            return 2
        return keep(args.deck, args.slide, args.brand_dir, args.slug, args.source)
    if args.command == "extract":
        if not args.deck.is_file():
            print(f"slide-library: no deck at {args.deck}", file=sys.stderr)
            return 2
        try:
            nummern = [int(teil) for teil in args.slides.split(",") if teil.strip()]
        except ValueError:
            print(f"slide-library: --slides takes numbers, got {args.slides!r}", file=sys.stderr)
            return 2
        return extract(args.deck, nummern, args.out)
    if args.command == "schema-check":
        if not args.deck.is_file():
            print(f"slide-library: no deck at {args.deck}", file=sys.stderr)
            return 2
        return schema_check(args.deck, args.slide)
    if args.command == "overflow-check":
        return overflow_check(args.deck, args.slide, args.baseline)
    if args.command == "overlap-check":
        return overlap_check(args.deck, args.slide, args.baseline)
    if args.command == "furniture-check":
        return furniture_check(args.deck, args.tol, args.zone)
    if args.command == "align-check":
        return align_check(args.deck, args.slide)
    if args.command == "render":
        if not args.deck.is_file():
            print(f"slide-library: {args.deck} is not a file", file=sys.stderr)
            return 1
        return render(args.deck, args.into, args.dpi, args.font_dirs)
    if args.command == "migrate":
        for pfad in [args.deck, args.into] + ([args.brand_from] if args.brand_from else []):
            if not pfad.is_file():
                print(f"slide-library: {pfad} is not a file", file=sys.stderr)
                return 1
        return migrate(args.deck, args.slide, args.into, args.out, args.layout, args.scheme,
                       args.title, args.replace, args.keep_shapes, args.brand_from)
    if args.command == "layout-check":
        return layout_check(args.deck, args.slide, args.min_margin, args.min_pt, args.min_pad)
    return show(args.library, args.slide)


if __name__ == "__main__":
    sys.exit(main(sys.argv))
